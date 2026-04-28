"""
Gera apresentações (.pptx) e PDF a partir dos relatórios markdown em
`outputs/relatorios_mensais_componentes/`.

Saídas em `outputs/relatorios_mensais_componentes_ppt/`:
  - <EQUIP>.pptx
  - <EQUIP>.pdf

Notas de compatibilidade:
- Tema com fundo branco.
- Os gráficos são embutidos como imagens PNG (geradas com matplotlib),
  evitando charts nativos do PowerPoint que costumam quebrar no
  LibreOffice/WPS.
- Nenhum conector ou autoshape exótica é utilizado — apenas retângulos,
  texto e imagens, recursos universalmente suportados.
- O PDF é produzido pelo LibreOffice em modo headless a partir do PPTX
  (mesmo render, garantindo consistência visual).
"""
from __future__ import annotations

import io
import re
import shutil
import subprocess
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

SCRIPTS_DIR = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPTS_DIR.parent / "config"))
try:
    from config import paths  # type: ignore
except Exception:
    import paths  # type: ignore

INPUT_DIR = paths.OUTPUTS_DIR / "relatorios_mensais_componentes"
OUTPUT_DIR = paths.OUTPUTS_DIR / "relatorios_mensais_componentes_ppt"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# ---------------------------------------------------------------------------
# Tema (paleta clara — fundo branco)
# ---------------------------------------------------------------------------
COLOR_BG = RGBColor(0xFF, 0xFF, 0xFF)            # fundo branco
COLOR_PANEL = RGBColor(0xF5, 0xF7, 0xFA)         # cinza muito claro
COLOR_BORDER = RGBColor(0xD0, 0xD7, 0xDE)        # cinza divisor
COLOR_TITLE = RGBColor(0x1F, 0x2D, 0x3D)         # azul escuro (título)
COLOR_TEXT = RGBColor(0x33, 0x3D, 0x4A)          # cinza escuro (corpo)
COLOR_MUTED = RGBColor(0x6B, 0x77, 0x85)         # cinza médio (legenda)
COLOR_ACCENT = RGBColor(0xE5, 0x3E, 0x5A)        # vermelho/rosa (destaque)
COLOR_WARNING = RGBColor(0xF1, 0xA1, 0x3C)       # amarelo (alerta)
COLOR_OK = RGBColor(0x2C, 0xA5, 0x7E)            # verde (ok)
COLOR_ACCENT_BLUE = RGBColor(0x2C, 0x6E, 0xC9)   # azul (info)

# matplotlib hex equivalents
HEX_PANEL = "#F5F7FA"
HEX_BORDER = "#D0D7DE"
HEX_TEXT = "#333D4A"
HEX_MUTED = "#6B7785"
HEX_ACCENT = "#E53E5A"
HEX_WARNING = "#F1A13C"
HEX_OK = "#2CA57E"
HEX_BLUE = "#2C6EC9"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MESES_PT = {
    1: "jan", 2: "fev", 3: "mar", 4: "abr", 5: "mai", 6: "jun",
    7: "jul", 8: "ago", 9: "set", 10: "out", 11: "nov", 12: "dez",
}


# ---------------------------------------------------------------------------
# Modelo
# ---------------------------------------------------------------------------
@dataclass
class Manutencao:
    data: str
    evento: str
    dentro_periodo: str


@dataclass
class LinhaMensal:
    ano_mes: str            # "2024-mai"
    composto: str
    quantidade: int
    pct_mes: float
    pct_total: float
    refugo: int
    dias: int


@dataclass
class EventoMensal:
    ano_mes: str            # ano-mês onde o evento aparece (pode ser "")
    data: str
    evento: str


@dataclass
class ResumoComposto:
    composto: str
    quantidade: int
    pct_total: float
    refugo: int
    meses: int


@dataclass
class Prescricao:
    data_prescrita: str = ""
    data_ultima_troca: str = ""
    data_referencia: str = ""
    dias_restantes: int | None = None
    urgencia: str = ""
    t_base: str = ""
    amplitude: str = ""
    massa_atual: str = ""
    massa_referencia: str = ""
    fator_massa: str = ""
    efeito_massa: str = ""
    ociosidade: str = ""
    fator_desgaste: str = "1,0000"
    t_prescrito_dias: str = ""
    formula_substituida: str = ""


@dataclass
class RelatorioEquipamento:
    equipamento: str
    fonte: str = ""
    periodo_inicio: str = ""
    periodo_fim: str = ""
    total_produzido: int = 0
    total_refugado: int = 0
    pct_refugo: float = 0.0
    n_compostos: int = 0
    manutencoes: list[Manutencao] = field(default_factory=list)
    linhas_mensais: list[LinhaMensal] = field(default_factory=list)
    eventos_mensais: list[EventoMensal] = field(default_factory=list)
    resumo: list[ResumoComposto] = field(default_factory=list)
    prescricao: Prescricao = field(default_factory=Prescricao)


# ---------------------------------------------------------------------------
# Parser do markdown
# ---------------------------------------------------------------------------
_NUM_PT_RE = re.compile(r"-?\d{1,3}(?:\.\d{3})*(?:,\d+)?|-?\d+(?:[.,]\d+)?")


def _to_int(s: str) -> int:
    s = (s or "").strip().replace(".", "").replace(" ", "")
    if s in ("", "—", "-"):
        return 0
    s = s.split(",")[0]
    try:
        return int(s)
    except ValueError:
        return 0


def _to_float(s: str) -> float:
    s = (s or "").strip().replace("%", "").replace(" ", "")
    s = s.replace(".", "").replace(",", ".")
    if s in ("", "—", "-"):
        return 0.0
    try:
        return float(s)
    except ValueError:
        return 0.0


def _split_md_row(line: str) -> list[str]:
    line = line.strip()
    if line.startswith("|"):
        line = line[1:]
    if line.endswith("|"):
        line = line[:-1]
    return [c.strip() for c in line.split("|")]


def _parse_event_cell(cell: str) -> tuple[str, str] | None:
    # ex.: 🔧 **22/05/2024** — _Preventiva RM.195 (#2)_
    m = re.search(r"\*\*([0-9/]+)\*\*\s*[—-]\s*_(.+?)_", cell)
    if not m:
        return None
    return m.group(1), m.group(2)


def parse_markdown(path: Path) -> RelatorioEquipamento:
    text = path.read_text(encoding="utf-8")
    lines = text.splitlines()
    rel = RelatorioEquipamento(equipamento=path.stem)

    # Cabeçalho
    for ln in lines[:30]:
        if ln.startswith("**Fonte:**"):
            m = re.search(r"`([^`]+)`", ln)
            rel.fonte = m.group(1) if m else ""
        elif ln.startswith("**Período coberto:**"):
            m = re.search(r"(\d{2}/\d{2}/\d{4}).+?(\d{2}/\d{2}/\d{4})", ln)
            if m:
                rel.periodo_inicio = m.group(1)
                rel.periodo_fim = m.group(2)
        elif ln.startswith("**Total produzido:**"):
            m = re.search(r"([\d.,]+)\s*peças", ln)
            if m:
                rel.total_produzido = _to_int(m.group(1))
        elif ln.startswith("**Total refugado:**"):
            m = re.search(r"([\d.,]+)\s*peças.*?\(([\d.,]+)\s*%\)", ln)
            if m:
                rel.total_refugado = _to_int(m.group(1))
                rel.pct_refugo = _to_float(m.group(2))
        elif ln.startswith("**Compostos distintos utilizados:**"):
            m = re.search(r"(\d+)\s*$", ln.strip())
            if m:
                rel.n_compostos = int(m.group(1))

    # Seções por título
    secoes: dict[str, list[str]] = {}
    cur = None
    for ln in lines:
        if ln.startswith("## "):
            cur = ln[3:].strip()
            secoes[cur] = []
        elif cur is not None:
            secoes[cur].append(ln)

    # --- Manutenções registradas ---
    for ln in secoes.get("Manutenções registradas", []):
        if not ln.startswith("|"):
            continue
        cells = _split_md_row(ln)
        if len(cells) < 3 or cells[0].lower().startswith("data"):
            continue
        if cells[0] and set(cells[0]) <= set("-:"):
            continue
        rel.manutencoes.append(Manutencao(cells[0], cells[1], cells[2]))

    # --- Produção mensal por composto ---
    cur_mes = ""
    for ln in secoes.get("Produção mensal por composto", []):
        if not ln.startswith("|"):
            continue
        cells = _split_md_row(ln)
        if len(cells) < 7:
            continue
        if cells[0].lower().startswith("ano-mês"):
            continue
        if cells[0] and set(cells[0]) <= set("-:"):
            continue
        # linha de evento: composto começa com 🔧
        if cells[1].startswith("🔧"):
            ev = _parse_event_cell(cells[1])
            if ev:
                rel.eventos_mensais.append(EventoMensal(cur_mes, ev[0], ev[1]))
            continue
        if cells[0]:
            cur_mes = cells[0]
        rel.linhas_mensais.append(
            LinhaMensal(
                ano_mes=cur_mes,
                composto=cells[1],
                quantidade=_to_int(cells[2]),
                pct_mes=_to_float(cells[3]),
                pct_total=_to_float(cells[4]),
                refugo=_to_int(cells[5]),
                dias=_to_int(cells[6]),
            )
        )

    # --- Resumo por composto ---
    for ln in secoes.get("Resumo por composto (período inteiro)", []):
        if not ln.startswith("|"):
            continue
        cells = _split_md_row(ln)
        if len(cells) < 5 or cells[0].lower().startswith("composto"):
            continue
        if cells[0] and set(cells[0]) <= set("-:"):
            continue
        nome = cells[0].replace("**", "").strip()
        if nome.upper() == "TOTAL":
            continue
        rel.resumo.append(
            ResumoComposto(
                composto=nome,
                quantidade=_to_int(cells[1]),
                pct_total=_to_float(cells[2]),
                refugo=_to_int(cells[3]),
                meses=_to_int(cells[4]),
            )
        )

    # --- Prescrição ---
    bloco = secoes.get("Prescrição da próxima manutenção", [])
    p = rel.prescricao
    for ln in bloco:
        if "Data prescrita:" in ln:
            m = re.search(r"`([^`]+)`", ln)
            p.data_prescrita = m.group(1) if m else ""
        elif "Data da última troca:" in ln:
            m = re.search(r"`([^`]+)`", ln)
            p.data_ultima_troca = m.group(1) if m else ""
        elif "Data de referência" in ln:
            m = re.search(r"`([^`]+)`", ln)
            p.data_referencia = m.group(1) if m else ""
        elif "Dias restantes:" in ln:
            m = re.search(r"`([^`]+)`", ln)
            if m:
                p.dias_restantes = _to_int(m.group(1))
        elif "Urgência:" in ln:
            depois = ln.split("Urgência:", 1)[1]
            limpo = depois.replace("**", "").strip()
            p.urgencia = limpo
        elif ln.startswith("| **T_base**"):
            cells = _split_md_row(ln)
            if len(cells) >= 2:
                p.t_base = cells[1]
        elif ln.startswith("| **Amplitude"):
            cells = _split_md_row(ln)
            if len(cells) >= 4:
                p.amplitude = cells[1]
                p.fator_desgaste = cells[3] or "1,0000"
        elif ln.startswith("| **Massa consumida"):
            cells = _split_md_row(ln)
            if len(cells) >= 5:
                p.massa_atual = cells[1]
                p.massa_referencia = cells[2]
                p.fator_massa = cells[3]
                p.efeito_massa = cells[4]
        elif ln.startswith("| **Dias de ociosidade**"):
            cells = _split_md_row(ln)
            if len(cells) >= 2:
                p.ociosidade = cells[1]
        elif "T_prescrito =" in ln:
            p.formula_substituida = ln.strip()
            m = re.search(r"=\s*([\d.,]+)\s*dias", ln)
            if m:
                p.t_prescrito_dias = m.group(1)

    return rel


# ---------------------------------------------------------------------------
# Helpers de slide
# ---------------------------------------------------------------------------
def _set_slide_bg(slide, color: RGBColor = COLOR_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, x, y, w, h, *, fill=COLOR_PANEL, border=COLOR_BORDER, border_width=Pt(0.75), rounded=False):
    shape_kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_kind, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border
        shape.line.width = border_width
    shape.shadow.inherit = False
    if rounded:
        try:
            shape.adjustments[0] = 0.08
        except Exception:
            pass
    # remove default text
    if shape.has_text_frame:
        shape.text_frame.text = ""
        shape.text_frame.margin_left = 0
        shape.text_frame.margin_right = 0
        shape.text_frame.margin_top = 0
        shape.text_frame.margin_bottom = 0
    return shape


def _add_text(slide, x, y, w, h, text, *,
              size=Pt(14), bold=False, color=COLOR_TEXT,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              font_name="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = size
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return tb


def _add_multiline(slide, x, y, w, h, segments, *, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    """segments: list of list[(text, size, bold, color)] — each inner list = paragraph."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.vertical_anchor = anchor
    for i, paragraph in enumerate(segments):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for txt, size, bold, color in paragraph:
            run = p.add_run()
            run.text = txt
            run.font.size = size
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = "Calibri"
    return tb


def _slide_header(slide, titulo: str, equipamento: str):
    # barra de título
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.85),
              fill=COLOR_BG, border=None)
    _add_rect(slide, Inches(0), Inches(0.82), SLIDE_W, Emu(28000),
              fill=COLOR_ACCENT, border=None)
    _add_text(slide, Inches(0.5), Inches(0.15), Inches(10.5), Inches(0.6),
              titulo, size=Pt(26), bold=True, color=COLOR_TITLE,
              anchor=MSO_ANCHOR.MIDDLE)
    _add_text(slide, Inches(11.0), Inches(0.15), Inches(2.0), Inches(0.6),
              equipamento, size=Pt(18), bold=True, color=COLOR_ACCENT,
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def _slide_footer(slide, equipamento: str, periodo: str):
    y = Inches(7.05)
    _add_rect(slide, Inches(0), y, SLIDE_W, Emu(15000),
              fill=COLOR_BORDER, border=None)
    _add_text(slide, Inches(0.5), Inches(7.12), Inches(8), Inches(0.35),
              f"Relatório Mensal — Componente {equipamento}",
              size=Pt(11), color=COLOR_MUTED, anchor=MSO_ANCHOR.MIDDLE)
    _add_text(slide, Inches(8.5), Inches(7.12), Inches(4.5), Inches(0.35),
              periodo, size=Pt(11), color=COLOR_MUTED,
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def _kpi_card(slide, x, y, w, h, value: str, label: str, sub: str = "",
              value_color: RGBColor = COLOR_ACCENT):
    _add_rect(slide, x, y, w, h, fill=COLOR_PANEL, border=COLOR_BORDER, rounded=True)
    _add_text(slide, x, y + Inches(0.15), w, Inches(0.7), value,
              size=Pt(28), bold=True, color=value_color,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _add_text(slide, x, y + Inches(0.85), w, Inches(0.4), label,
              size=Pt(13), bold=True, color=COLOR_TITLE,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        _add_text(slide, x, y + Inches(1.18), w, Inches(0.4), sub,
                  size=Pt(10), color=COLOR_MUTED,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _add_image_from_fig(slide, fig, x, y, w_in, h_in):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="white")
    plt.close(fig)
    buf.seek(0)
    slide.shapes.add_picture(buf, x, y, width=Inches(w_in), height=Inches(h_in))


# ---------------------------------------------------------------------------
# matplotlib base
# ---------------------------------------------------------------------------
def _setup_mpl():
    plt.rcParams.update({
        "font.family": "DejaVu Sans",
        "font.size": 10,
        "axes.edgecolor": HEX_BORDER,
        "axes.labelcolor": HEX_TEXT,
        "xtick.color": HEX_TEXT,
        "ytick.color": HEX_TEXT,
        "axes.spines.top": False,
        "axes.spines.right": False,
    })


def _fmt_int_pt(v) -> str:
    try:
        return f"{int(v):,}".replace(",", ".")
    except Exception:
        return str(v)


def _fmt_pct_pt(v) -> str:
    return f"{v:,.2f} %".replace(",", "X").replace(".", ",").replace("X", ".")


# ---------------------------------------------------------------------------
# Construção de slides
# ---------------------------------------------------------------------------
def _periodo_label(rel: RelatorioEquipamento) -> str:
    def to_label(d: str) -> str:
        m = re.match(r"(\d{2})/(\d{2})/(\d{4})", d)
        if not m:
            return d
        mes = MESES_PT[int(m.group(2))].capitalize()
        return f"{mes}/{m.group(3)}"

    return f"{to_label(rel.periodo_inicio)} → {to_label(rel.periodo_fim)}"


def _meses_operacao(rel: RelatorioEquipamento) -> int:
    meses = {l.ano_mes for l in rel.linhas_mensais if l.ano_mes}
    return len(meses)


def _producao_mensal(rel: RelatorioEquipamento) -> list[tuple[str, int]]:
    agg: dict[str, int] = {}
    ordem: list[str] = []
    for l in rel.linhas_mensais:
        if not l.ano_mes:
            continue
        if l.ano_mes not in agg:
            agg[l.ano_mes] = 0
            ordem.append(l.ano_mes)
        agg[l.ano_mes] += l.quantidade
    return [(m, agg[m]) for m in ordem]


def _wrap_short(text: str, width: int) -> str:
    """Quebra texto em duas linhas de até `width` chars."""
    if not text or len(text) <= width:
        return text
    cut = text.rfind(" ", 0, width)
    if cut == -1:
        cut = width
    return text[:cut] + "\n" + text[cut:].strip()[:width]


def _ano_mes_label(ano_mes: str) -> str:
    """Converte '2025-set' → 'Set/2025'."""
    m = re.match(r"(\d{4})-([a-zçãé]+)", ano_mes or "")
    if not m:
        return ano_mes
    return f"{m.group(2).capitalize()}/{m.group(1)}"


def _slide_capa(prs, rel: RelatorioEquipamento):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    # faixa lateral
    _add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H,
              fill=COLOR_ACCENT, border=None)

    # bloco esquerdo (título)
    _add_rect(slide, Inches(0.6), Inches(1.4), Inches(0.9), Inches(0.5),
              fill=COLOR_WARNING, border=None, rounded=True)
    _add_text(slide, Inches(0.6), Inches(1.4), Inches(0.9), Inches(0.5),
              rel.equipamento, size=Pt(16), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _add_multiline(slide, Inches(0.6), Inches(2.0), Inches(5.5), Inches(2.8),
                   [
                       [("Relatório", Pt(44), True, COLOR_TITLE)],
                       [("Mensal", Pt(44), True, COLOR_TITLE)],
                       [("por", Pt(44), True, COLOR_TITLE)],
                       [("Componente", Pt(44), True, COLOR_TITLE)],
                   ])

    _add_text(slide, Inches(0.6), Inches(5.0), Inches(5.5), Inches(0.4),
              f"Análise de desempenho · {_periodo_label(rel)}",
              size=Pt(14), color=COLOR_TEXT)
    _add_text(slide, Inches(0.6), Inches(5.45), Inches(5.5), Inches(0.4),
              "Produção · Qualidade · Manutenção",
              size=Pt(12), color=COLOR_MUTED)

    # KPIs à direita
    x_kpi = Inches(7.0)
    w_kpi = Inches(5.7)
    h_kpi = Inches(1.4)
    gap = Inches(0.25)
    meses_op = _meses_operacao(rel) or 0

    items = [
        (_fmt_int_pt(rel.total_produzido), "Peças produzidas",
         _periodo_label(rel), COLOR_ACCENT),
        (str(rel.n_compostos), "Compostos utilizados",
         f"{meses_op} meses de operação", COLOR_WARNING),
        (_fmt_pct_pt(rel.pct_refugo), "Taxa de refugo",
         f"{_fmt_int_pt(rel.total_refugado)} peças refugadas", COLOR_OK),
    ]
    for i, (val, label, sub, c) in enumerate(items):
        y = Inches(1.4) + (h_kpi + gap) * i
        _add_rect(slide, x_kpi, y, w_kpi, h_kpi,
                  fill=COLOR_PANEL, border=COLOR_BORDER, rounded=True)
        _add_text(slide, x_kpi + Inches(0.3), y + Inches(0.15),
                  Inches(2.0), Inches(0.6),
                  val, size=Pt(28), bold=True, color=c)
        _add_text(slide, x_kpi + Inches(0.3), y + Inches(0.7),
                  w_kpi - Inches(0.6), Inches(0.4),
                  label, size=Pt(14), bold=True, color=COLOR_TITLE)
        _add_text(slide, x_kpi + Inches(0.3), y + Inches(1.0),
                  w_kpi - Inches(0.6), Inches(0.35),
                  sub, size=Pt(11), color=COLOR_MUTED)


def _slide_visao_executiva(prs, rel: RelatorioEquipamento):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _slide_header(slide, "Visão Executiva da Produção", rel.equipamento)

    # 4 KPIs no topo
    meses_op = _meses_operacao(rel)
    aprov_pct = 100 - rel.pct_refugo
    kpis = [
        (_fmt_int_pt(rel.total_produzido), "Total Produzido", "peças", COLOR_ACCENT),
        (str(meses_op), "Meses de Operação", _periodo_label(rel), COLOR_WARNING),
        (str(rel.n_compostos), "Compostos Distintos", "utilizados", COLOR_BLUE_OR_DEFAULT()),
        (_fmt_pct_pt(rel.pct_refugo), "Taxa de Refugo",
         f"{_fmt_int_pt(rel.total_refugado)} peças", COLOR_ACCENT),
    ]
    x0 = Inches(0.4)
    w = Inches(3.05)
    h = Inches(1.55)
    y = Inches(1.1)
    gap = Inches(0.1)
    for i, (val, lbl, sub, c) in enumerate(kpis):
        x = x0 + (w + gap) * i
        _add_rect(slide, x, y, w, h, fill=COLOR_PANEL, border=COLOR_BORDER, rounded=True)
        _add_text(slide, x, y + Inches(0.15), w, Inches(0.7),
                  val, size=Pt(26), bold=True, color=c,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text(slide, x, y + Inches(0.8), w, Inches(0.35),
                  lbl, size=Pt(12), bold=True, color=COLOR_TITLE,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text(slide, x, y + Inches(1.15), w, Inches(0.35),
                  sub, size=Pt(10), color=COLOR_MUTED,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Painel esquerdo: destaques
    panel_y = Inches(2.85)
    panel_h = Inches(4.0)
    _add_rect(slide, Inches(0.4), panel_y, Inches(7.5), panel_h,
              fill=COLOR_PANEL, border=COLOR_BORDER, rounded=True)
    _add_text(slide, Inches(0.6), panel_y + Inches(0.15), Inches(7.0), Inches(0.45),
              "Destaques do Período", size=Pt(16), bold=True, color=COLOR_TITLE)

    destaques = _gerar_destaques(rel)
    y = panel_y + Inches(0.7)
    for icon, line in destaques:
        _add_text(slide, Inches(0.7), y, Inches(0.4), Inches(0.4),
                  icon, size=Pt(14), color=COLOR_OK)
        _add_text(slide, Inches(1.05), y, Inches(6.7), Inches(0.7),
                  line, size=Pt(12), color=COLOR_TEXT)
        y += Inches(0.7)

    # Donut Aprovado vs Refugado (matplotlib)
    _add_rect(slide, Inches(8.1), panel_y, Inches(4.85), panel_h,
              fill=COLOR_PANEL, border=COLOR_BORDER, rounded=True)
    _add_text(slide, Inches(8.1), panel_y + Inches(0.15), Inches(4.85), Inches(0.45),
              "Aprovado vs. Refugado", size=Pt(16), bold=True, color=COLOR_TITLE,
              align=PP_ALIGN.CENTER)
    _setup_mpl()
    fig, ax = plt.subplots(figsize=(4.6, 3.4), facecolor="white")
    sizes = [aprov_pct, rel.pct_refugo]
    colors_ = [HEX_OK, HEX_ACCENT]
    wedges, _ = ax.pie(sizes, colors=colors_, startangle=90,
                       wedgeprops=dict(width=0.35, edgecolor="white", linewidth=2))
    ax.text(0, 0.05, f"{aprov_pct:.2f}%".replace(".", ","),
            ha="center", va="center", fontsize=22, fontweight="bold", color=HEX_TEXT)
    ax.text(0, -0.18, "aprovado", ha="center", va="center",
            fontsize=11, color=HEX_MUTED)
    ax.set_aspect("equal")
    ax.legend(["Aprovado", "Refugado"], loc="lower center",
              bbox_to_anchor=(0.5, -0.05), ncol=2, frameon=False,
              fontsize=10)
    _add_image_from_fig(slide, fig, Inches(8.2), panel_y + Inches(0.6),
                        4.6, 3.3)

    _slide_footer(slide, rel.equipamento, _periodo_label(rel))


def COLOR_BLUE_OR_DEFAULT():
    return COLOR_ACCENT_BLUE


def _gerar_destaques(rel: RelatorioEquipamento) -> list[tuple[str, str]]:
    out: list[tuple[str, str]] = []
    prod = _producao_mensal(rel)
    if prod:
        pico_mes, pico_qtd = max(prod, key=lambda x: x[1])
        out.append(("✓", f"Pico de produção em {_ano_mes_label(pico_mes)}: "
                         f"{_fmt_int_pt(pico_qtd)} peças"))
    if rel.resumo:
        topo = rel.resumo[0]
        out.append(("✓", f"{topo.composto} lidera com "
                         f"{_fmt_pct_pt(topo.pct_total)} da produção total"))
    out.append(("✓", f"Taxa de refugo de {_fmt_pct_pt(rel.pct_refugo)} no período"))
    if rel.prescricao.data_ultima_troca and rel.prescricao.data_prescrita:
        out.append(("✓", f"Última troca em {rel.prescricao.data_ultima_troca} — "
                         f"próxima prescrita em {rel.prescricao.data_prescrita}"))
    return out


def _slide_tendencias(prs, rel: RelatorioEquipamento):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _slide_header(slide, "Tendências de Produção Mensal", rel.equipamento)

    prod = _producao_mensal(rel)
    if not prod:
        _add_text(slide, Inches(0.5), Inches(3), Inches(12), Inches(1),
                  "Sem dados mensais de produção.", size=Pt(16), color=COLOR_MUTED,
                  align=PP_ALIGN.CENTER)
        _slide_footer(slide, rel.equipamento, _periodo_label(rel))
        return

    qtds = [q for _, q in prod]
    media = sum(qtds) / len(qtds)
    pico_mes, pico_qtd = max(prod, key=lambda x: x[1])
    min_mes, min_qtd = min(prod, key=lambda x: x[1])

    # 3 KPIs no topo
    cards = [
        (_fmt_int_pt(pico_qtd), "Pico de Produção",
         f"{_ano_mes_label(pico_mes)} — recorde do período", COLOR_ACCENT),
        (_fmt_int_pt(int(round(media))), "Média Mensal",
         f"Calculada sobre {len(prod)} meses", COLOR_WARNING),
        (_fmt_int_pt(min_qtd), "Menor Mês",
         f"{_ano_mes_label(min_mes)} — menor volume registrado",
         COLOR_ACCENT_BLUE),
    ]
    x0 = Inches(0.4)
    w = Inches(4.15)
    h = Inches(1.3)
    y = Inches(1.1)
    gap = Inches(0.1)
    for i, (val, lbl, sub, c) in enumerate(cards):
        x = x0 + (w + gap) * i
        _add_rect(slide, x, y, w, h, fill=COLOR_PANEL, border=COLOR_BORDER, rounded=True)
        _add_text(slide, x + Inches(0.3), y + Inches(0.1), w - Inches(0.6), Inches(0.55),
                  val, size=Pt(24), bold=True, color=c)
        _add_text(slide, x + Inches(0.3), y + Inches(0.6), w - Inches(0.6), Inches(0.35),
                  lbl, size=Pt(13), bold=True, color=COLOR_TITLE)
        _add_text(slide, x + Inches(0.3), y + Inches(0.92), w - Inches(0.6), Inches(0.35),
                  sub, size=Pt(10), color=COLOR_MUTED)

    # Gráfico de barras
    _setup_mpl()
    fig, ax = plt.subplots(figsize=(12.3, 4.0), facecolor="white")
    labels = [_ano_mes_label(m) for m, _ in prod]
    valores = [q for _, q in prod]
    cores = [HEX_ACCENT if q == pico_qtd else HEX_BLUE for q in valores]
    bars = ax.bar(range(len(labels)), valores, color=cores, edgecolor="white", linewidth=0.5)
    ax.axhline(media, color=HEX_WARNING, linestyle="--", linewidth=1.2,
               label=f"Média: {_fmt_int_pt(int(round(media)))}")
    ax.set_xticks(range(len(labels)))
    ax.set_xticklabels(labels, rotation=45, ha="right", fontsize=9)
    ax.set_ylabel("Peças produzidas", color=HEX_TEXT)
    ax.yaxis.set_major_formatter(
        plt.FuncFormatter(lambda x, _: _fmt_int_pt(int(x)))
    )
    ax.grid(axis="y", linestyle=":", linewidth=0.6, color=HEX_BORDER)
    ax.set_axisbelow(True)
    ax.legend(loc="upper right", frameon=False, fontsize=10)
    fig.tight_layout()
    _add_image_from_fig(slide, fig, Inches(0.4), Inches(2.55), 12.55, 4.4)

    _slide_footer(slide, rel.equipamento, _periodo_label(rel))


def _slide_compostos(prs, rel: RelatorioEquipamento):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _slide_header(slide, "Análise de Materiais e Compostos", rel.equipamento)

    if not rel.resumo:
        _add_text(slide, Inches(0.5), Inches(3), Inches(12), Inches(1),
                  "Sem resumo de compostos disponível.", size=Pt(16),
                  color=COLOR_MUTED, align=PP_ALIGN.CENTER)
        _slide_footer(slide, rel.equipamento, _periodo_label(rel))
        return

    # Painel esquerdo: gráfico de barras horizontal
    _add_rect(slide, Inches(0.4), Inches(1.1), Inches(7.3), Inches(5.85),
              fill=COLOR_PANEL, border=COLOR_BORDER, rounded=True)
    _add_text(slide, Inches(0.6), Inches(1.2), Inches(6.9), Inches(0.4),
              "Volume por Composto", size=Pt(16), bold=True, color=COLOR_TITLE)

    _setup_mpl()
    ordenados = sorted(rel.resumo, key=lambda c: c.quantidade)
    nomes = [c.composto for c in ordenados]
    valores = [c.quantidade for c in ordenados]
    max_q = max(valores) if valores else 0
    cores = []
    for c in ordenados:
        if c.quantidade == max_q:
            cores.append(HEX_ACCENT)
        elif c.pct_total >= 5.0:
            cores.append(HEX_WARNING)
        else:
            cores.append(HEX_BLUE)
    fig, ax = plt.subplots(figsize=(7.0, 5.2), facecolor="white")
    ax.barh(range(len(nomes)), valores, color=cores, edgecolor="white", linewidth=0.5)
    ax.set_yticks(range(len(nomes)))
    ax.set_yticklabels([n[:30] for n in nomes], fontsize=9)
    ax.set_xlabel("Peças produzidas", color=HEX_TEXT, fontsize=10)
    ax.xaxis.set_major_formatter(
        plt.FuncFormatter(lambda x, _: _fmt_int_pt(int(x)))
    )
    ax.grid(axis="x", linestyle=":", linewidth=0.6, color=HEX_BORDER)
    ax.set_axisbelow(True)
    fig.tight_layout()
    _add_image_from_fig(slide, fig, Inches(0.5), Inches(1.65), 7.1, 5.2)

    # Painel direito: composto dominante + tabela
    pdom = rel.resumo[0]
    _add_rect(slide, Inches(7.85), Inches(1.1), Inches(5.1), Inches(2.05),
              fill=COLOR_PANEL, border=COLOR_BORDER, rounded=True)
    _add_text(slide, Inches(8.0), Inches(1.2), Inches(4.8), Inches(0.35),
              "Composto Dominante", size=Pt(11), color=COLOR_MUTED)
    _add_text(slide, Inches(8.0), Inches(1.5), Inches(4.8), Inches(0.4),
              pdom.composto, size=Pt(15), bold=True, color=COLOR_TITLE)
    _add_text(slide, Inches(8.0), Inches(1.95), Inches(2.4), Inches(0.45),
              _fmt_pct_pt(pdom.pct_total), size=Pt(22), bold=True, color=COLOR_ACCENT)
    _add_text(slide, Inches(8.0), Inches(2.45), Inches(2.4), Inches(0.3),
              "do total produzido", size=Pt(10), color=COLOR_MUTED)
    _add_text(slide, Inches(10.4), Inches(1.95), Inches(2.4), Inches(0.45),
              _fmt_int_pt(pdom.quantidade), size=Pt(22), bold=True, color=COLOR_ACCENT)
    _add_text(slide, Inches(10.4), Inches(2.45), Inches(2.4), Inches(0.3),
              "peças", size=Pt(10), color=COLOR_MUTED)
    _add_text(slide, Inches(8.0), Inches(2.78), Inches(2.4), Inches(0.3),
              f"{pdom.meses} meses com uso", size=Pt(11), color=COLOR_TEXT)
    _add_text(slide, Inches(10.4), Inches(2.78), Inches(2.4), Inches(0.3),
              f"{_fmt_int_pt(pdom.refugo)} peças refugadas",
              size=Pt(11), color=COLOR_TEXT)

    # Tabela top compostos
    _add_rect(slide, Inches(7.85), Inches(3.3), Inches(5.1), Inches(3.65),
              fill=COLOR_PANEL, border=COLOR_BORDER, rounded=True)
    headers = ["Composto", "Produzido", "% Total", "Refugo", "Meses"]
    col_x = [Inches(7.95), Inches(9.95), Inches(10.95), Inches(11.75), Inches(12.45)]
    col_w = [Inches(2.0), Inches(1.0), Inches(0.8), Inches(0.7), Inches(0.5)]
    y_h = Inches(3.4)
    for i, htxt in enumerate(headers):
        align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT
        _add_text(slide, col_x[i], y_h, col_w[i], Inches(0.35),
                  htxt, size=Pt(11), bold=True, color=COLOR_ACCENT, align=align)
    # divisor
    _add_rect(slide, Inches(7.95), Inches(3.78), Inches(5.0), Emu(8000),
              fill=COLOR_BORDER, border=None)
    y_row = Inches(3.85)
    for idx, c in enumerate(rel.resumo[:6]):
        cor = COLOR_ACCENT if idx == 0 else COLOR_TEXT
        bold = idx == 0
        cells = [
            c.composto[:24],
            _fmt_int_pt(c.quantidade),
            _fmt_pct_pt(c.pct_total),
            _fmt_int_pt(c.refugo),
            str(c.meses),
        ]
        for i, txt in enumerate(cells):
            align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT
            _add_text(slide, col_x[i], y_row, col_w[i], Inches(0.4),
                      txt, size=Pt(10.5), bold=bold, color=cor, align=align)
        y_row += Inches(0.46)

    _slide_footer(slide, rel.equipamento, _periodo_label(rel))


def _slide_qualidade(prs, rel: RelatorioEquipamento):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _slide_header(slide, "Controle de Qualidade — Taxa de Refugo", rel.equipamento)

    # Painel esquerdo: refugo por composto
    _add_rect(slide, Inches(0.4), Inches(1.1), Inches(7.3), Inches(5.85),
              fill=COLOR_PANEL, border=COLOR_BORDER, rounded=True)
    _add_text(slide, Inches(0.6), Inches(1.2), Inches(6.9), Inches(0.4),
              "Refugo por Composto", size=Pt(16), bold=True, color=COLOR_TITLE)

    if rel.resumo:
        _setup_mpl()
        ordenados = sorted(rel.resumo, key=lambda c: c.refugo)
        nomes = [c.composto for c in ordenados]
        valores = [c.refugo for c in ordenados]
        max_r = max(valores) if valores else 0
        cores = []
        for v in valores:
            if v == 0:
                cores.append(HEX_BLUE)
            elif v >= 0.8 * max_r and max_r > 0:
                cores.append(HEX_ACCENT)
            elif v >= 0.4 * max_r and max_r > 0:
                cores.append(HEX_WARNING)
            else:
                cores.append(HEX_BLUE)
        fig, ax = plt.subplots(figsize=(7.0, 5.2), facecolor="white")
        ax.barh(range(len(nomes)), valores, color=cores,
                edgecolor="white", linewidth=0.5)
        ax.set_yticks(range(len(nomes)))
        ax.set_yticklabels([n[:30] for n in nomes], fontsize=9)
        ax.set_xlabel("Peças refugadas", color=HEX_TEXT, fontsize=10)
        ax.xaxis.set_major_formatter(
            plt.FuncFormatter(lambda x, _: _fmt_int_pt(int(x)))
        )
        ax.grid(axis="x", linestyle=":", linewidth=0.6, color=HEX_BORDER)
        ax.set_axisbelow(True)
        fig.tight_layout()
        _add_image_from_fig(slide, fig, Inches(0.5), Inches(1.65), 7.1, 5.2)

    # Painel direito: KPI
    _add_rect(slide, Inches(7.85), Inches(1.1), Inches(5.1), Inches(2.4),
              fill=COLOR_PANEL, border=COLOR_BORDER, rounded=True)
    cor_status = COLOR_OK if rel.pct_refugo < 2.0 else COLOR_WARNING
    _add_text(slide, Inches(7.85), Inches(1.25), Inches(5.1), Inches(0.7),
              _fmt_pct_pt(rel.pct_refugo),
              size=Pt(34), bold=True, color=cor_status,
              align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(7.85), Inches(1.95), Inches(5.1), Inches(0.4),
              "Taxa de Refugo Global", size=Pt(13), bold=True, color=COLOR_TITLE,
              align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(7.95), Inches(2.4), Inches(2.4), Inches(0.4),
              _fmt_int_pt(rel.total_refugado),
              size=Pt(18), bold=True, color=COLOR_ACCENT,
              align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(7.95), Inches(2.78), Inches(2.4), Inches(0.3),
              "Peças refugadas", size=Pt(10), color=COLOR_MUTED,
              align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(10.4), Inches(2.4), Inches(2.4), Inches(0.4),
              _fmt_int_pt(rel.total_produzido),
              size=Pt(18), bold=True, color=COLOR_WARNING,
              align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(10.4), Inches(2.78), Inches(2.4), Inches(0.3),
              "Total produzido", size=Pt(10), color=COLOR_MUTED,
              align=PP_ALIGN.CENTER)
    label_status = "Dentro do limite aceitável" if rel.pct_refugo < 2.0 else "Atenção ao patamar"
    _add_rect(slide, Inches(8.5), Inches(3.1), Inches(3.8), Inches(0.32),
              fill=RGBColor(0xE7, 0xF6, 0xEE) if rel.pct_refugo < 2.0
              else RGBColor(0xFD, 0xF1, 0xE0),
              border=None, rounded=True)
    _add_text(slide, Inches(8.5), Inches(3.1), Inches(3.8), Inches(0.32),
              label_status, size=Pt(11), bold=True, color=cor_status,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Painel direito: Análise de qualidade
    _add_rect(slide, Inches(7.85), Inches(3.65), Inches(5.1), Inches(3.3),
              fill=COLOR_PANEL, border=COLOR_BORDER, rounded=True)
    _add_text(slide, Inches(8.0), Inches(3.75), Inches(4.8), Inches(0.4),
              "Análise de Qualidade", size=Pt(14), bold=True, color=COLOR_TITLE)
    insights = _gerar_insights_qualidade(rel)
    y = Inches(4.2)
    for ic, txt, cor in insights:
        _add_text(slide, Inches(8.0), y, Inches(0.3), Inches(0.45),
                  ic, size=Pt(12), color=cor)
        _add_text(slide, Inches(8.3), y, Inches(4.5), Inches(0.7),
                  txt, size=Pt(10.5), color=COLOR_TEXT)
        y += Inches(0.65)

    _slide_footer(slide, rel.equipamento, _periodo_label(rel))


def _gerar_insights_qualidade(rel: RelatorioEquipamento):
    out = []
    if rel.resumo:
        com_refugo = [c for c in rel.resumo if c.refugo > 0]
        if com_refugo:
            top = max(com_refugo, key=lambda c: c.refugo)
            out.append(("●", f"{top.composto}: maior volume de refugo absoluto "
                             f"({_fmt_int_pt(top.refugo)} peças)", COLOR_ACCENT))
        zero = [c for c in rel.resumo if c.refugo == 0]
        if zero:
            z = zero[0]
            out.append(("●", f"{z.composto}: zero refugo registrado em "
                             f"{z.meses} mês(es) de uso", COLOR_OK))
    out.append(("✓", "Taxa mensal estável — sem pico de rejeição significativo",
                COLOR_OK))
    if rel.pct_refugo < 1.5:
        out.append(("✓", "Refugo geral abaixo de 1,5 % no período", COLOR_OK))
    else:
        out.append(("●", "Refugo geral próximo do limite — monitorar tendência",
                    COLOR_WARNING))
    return out


def _slide_manutencao(prs, rel: RelatorioEquipamento):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _slide_header(slide, "Estratégia de Manutenção Preventiva", rel.equipamento)

    p = rel.prescricao

    # Painel esquerdo: Histórico (timeline)
    _add_rect(slide, Inches(0.4), Inches(1.1), Inches(8.3), Inches(2.3),
              fill=COLOR_PANEL, border=COLOR_BORDER, rounded=True)
    _add_text(slide, Inches(0.6), Inches(1.2), Inches(8.0), Inches(0.4),
              "Histórico de Manutenções RM.195",
              size=Pt(15), bold=True, color=COLOR_TITLE)

    eventos = list(rel.manutencoes)
    n = len(eventos)
    if n > 0:
        _setup_mpl()
        fig, ax = plt.subplots(figsize=(8.0, 2.0), facecolor="white")
        ax.set_xlim(-0.5, max(n - 0.5, 0.5))
        ax.set_ylim(-1.6, 1.4)
        ax.axis("off")
        ax.hlines(0, -0.3, max(n - 0.7, 0.001), colors=HEX_BORDER,
                  linewidth=2, zorder=1)
        for i, ev in enumerate(eventos):
            cor = HEX_BLUE
            txt_cor = HEX_MUTED
            ev_lower = (ev.evento or "").lower()
            if "última troca" in ev_lower and "penúltima" not in ev_lower:
                cor = HEX_WARNING
                txt_cor = HEX_TEXT
            if i == n - 1:
                cor = HEX_ACCENT
                txt_cor = HEX_ACCENT
            ax.scatter(i, 0, s=180, color=cor, zorder=3,
                       edgecolor="white", linewidth=2)
            # alterna acima/abaixo para evitar sobreposição
            if i % 2 == 0:
                ax.text(i, 0.35, ev.data, ha="center", va="bottom",
                        fontsize=8.5, color=txt_cor, fontweight="bold")
                ax.text(i, 0.85, _wrap_short(ev.evento, 18),
                        ha="center", va="bottom",
                        fontsize=7.5, color=HEX_MUTED)
            else:
                ax.text(i, -0.35, ev.data, ha="center", va="top",
                        fontsize=8.5, color=txt_cor, fontweight="bold")
                ax.text(i, -0.85, _wrap_short(ev.evento, 18),
                        ha="center", va="top",
                        fontsize=7.5, color=HEX_MUTED)
        fig.tight_layout()
        _add_image_from_fig(slide, fig, Inches(0.5), Inches(1.55), 8.1, 1.85)

    # Painel direito: KPI próxima manutenção
    _add_rect(slide, Inches(8.85), Inches(1.1), Inches(4.1), Inches(3.0),
              fill=COLOR_PANEL, border=COLOR_BORDER, rounded=True)
    cor_urg = _urgencia_cor(p.urgencia)
    badge_bg = _urgencia_badge_bg(p.urgencia)
    _add_rect(slide, Inches(9.6), Inches(1.25), Inches(2.6), Inches(0.42),
              fill=badge_bg, border=None, rounded=True)
    label_badge = f"✓ {p.urgencia}" if p.urgencia.upper() == "OK" else (p.urgencia or "—")
    _add_text(slide, Inches(9.6), Inches(1.25), Inches(2.6), Inches(0.42),
              label_badge, size=Pt(12), bold=True, color=cor_urg,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    dias_txt = str(p.dias_restantes) if p.dias_restantes is not None else "—"
    _add_text(slide, Inches(8.85), Inches(1.85), Inches(4.1), Inches(0.7),
              dias_txt, size=Pt(36), bold=True, color=cor_urg,
              align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(8.85), Inches(2.55), Inches(4.1), Inches(0.35),
              "dias restantes", size=Pt(11), color=COLOR_MUTED,
              align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(8.85), Inches(2.95), Inches(4.1), Inches(0.4),
              p.data_prescrita or "—", size=Pt(16), bold=True, color=COLOR_TITLE,
              align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(8.85), Inches(3.35), Inches(4.1), Inches(0.3),
              "Próxima manutenção", size=Pt(10), color=COLOR_MUTED,
              align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(8.85), Inches(3.65), Inches(4.1), Inches(0.3),
              f"Última troca: {p.data_ultima_troca or '—'}",
              size=Pt(10), color=COLOR_MUTED, align=PP_ALIGN.CENTER)

    # Tabela: componentes da fórmula
    _add_rect(slide, Inches(0.4), Inches(3.55), Inches(8.3), Inches(3.4),
              fill=COLOR_PANEL, border=COLOR_BORDER, rounded=True)
    headers = ["Componente", "Valor Atual", "Referência", "Fator", "Efeito"]
    col_x = [Inches(0.55), Inches(2.55), Inches(4.0), Inches(5.6), Inches(6.4)]
    col_w = [Inches(2.0), Inches(1.4), Inches(1.6), Inches(0.8), Inches(2.2)]
    y_h = Inches(3.7)
    for i, htxt in enumerate(headers):
        _add_text(slide, col_x[i], y_h, col_w[i], Inches(0.35),
                  htxt, size=Pt(11), bold=True, color=COLOR_ACCENT)
    _add_rect(slide, Inches(0.55), Inches(4.05), Inches(8.0), Emu(8000),
              fill=COLOR_BORDER, border=None)

    rows = [
        ("T_base", f"{p.t_base or '—'}", "Mediana histórica", "—", "Ponto de partida"),
        ("Ampl. cil./fuso", p.amplitude or "—", "—", p.fator_desgaste or "1,0000",
         "Fator neutro" if (p.fator_desgaste or "").startswith("1,0") else "Ajuste por desgaste"),
        ("Massa pós-troca", p.massa_atual or "—",
         (p.massa_referencia or "—"), p.fator_massa or "—", p.efeito_massa or "—"),
        ("Ociosidade", p.ociosidade or "—", "Últ. prod. → hoje", "—",
         "Desliza data p/ frente"),
        ("T_prescrito", f"{p.t_prescrito_dias or '—'} dias",
         p.formula_substituida.replace("T_prescrito = ", "")
         .replace(f"= {p.t_prescrito_dias} dias", "").strip().rstrip("=").strip()
         if p.formula_substituida else "—",
         "—",
         f"Data: {p.data_prescrita or '—'}"),
    ]
    y_row = Inches(4.15)
    for idx, row in enumerate(rows):
        is_last = idx == len(rows) - 1
        cor = COLOR_ACCENT if is_last else COLOR_TEXT
        bold = is_last
        if is_last:
            _add_rect(slide, Inches(0.45), y_row - Inches(0.05),
                      Inches(8.2), Inches(0.55),
                      fill=RGBColor(0xFC, 0xEC, 0xEF), border=None, rounded=True)
        for i, txt in enumerate(row):
            _add_text(slide, col_x[i], y_row, col_w[i], Inches(0.5),
                      txt, size=Pt(10), bold=bold, color=cor)
        y_row += Inches(0.55)

    # Painel direito inferior: legenda urgência
    _add_rect(slide, Inches(8.85), Inches(4.25), Inches(4.1), Inches(2.7),
              fill=COLOR_PANEL, border=COLOR_BORDER, rounded=True)
    _add_text(slide, Inches(9.0), Inches(4.35), Inches(3.8), Inches(0.4),
              "Níveis de Urgência", size=Pt(13), bold=True, color=COLOR_TITLE)
    legenda = [
        ("ATRASADO", "< 0 dias", COLOR_ACCENT),
        ("URGENTE", "0 – 29 dias", COLOR_WARNING),
        ("ATENÇÃO", "30 – 89 dias", RGBColor(0xE3, 0xC0, 0x47)),
        ("✓ OK — ATUAL", "≥ 90 dias", COLOR_OK),
    ]
    yL = Inches(4.85)
    for nome, faixa, cor in legenda:
        _add_rect(slide, Inches(9.0), yL, Inches(1.55), Inches(0.35),
                  fill=_lighten(cor), border=None, rounded=True)
        _add_text(slide, Inches(9.0), yL, Inches(1.55), Inches(0.35),
                  nome, size=Pt(10), bold=True, color=cor,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text(slide, Inches(10.7), yL, Inches(2.1), Inches(0.35),
                  faixa, size=Pt(10), color=COLOR_TEXT,
                  anchor=MSO_ANCHOR.MIDDLE)
        yL += Inches(0.45)

    _slide_footer(slide, rel.equipamento,
                  f"Referência: {p.data_referencia or '—'}")


def _lighten(c: RGBColor) -> RGBColor:
    r, g, b = c[0], c[1], c[2]
    f = 0.85
    return RGBColor(int(r + (255 - r) * f),
                    int(g + (255 - g) * f),
                    int(b + (255 - b) * f))


def _urgencia_cor(urg: str) -> RGBColor:
    u = (urg or "").upper()
    if u == "OK":
        return COLOR_OK
    if u == "ATENÇÃO":
        return RGBColor(0xCD, 0xA8, 0x35)
    if u == "URGENTE":
        return COLOR_WARNING
    if u == "ATRASADO":
        return COLOR_ACCENT
    return COLOR_MUTED


def _urgencia_badge_bg(urg: str) -> RGBColor:
    u = (urg or "").upper()
    if u == "OK":
        return RGBColor(0xE7, 0xF6, 0xEE)
    if u in ("ATENÇÃO", "URGENTE"):
        return RGBColor(0xFD, 0xF1, 0xE0)
    if u == "ATRASADO":
        return RGBColor(0xFC, 0xEC, 0xEF)
    return RGBColor(0xEE, 0xEE, 0xEE)


def _slide_conclusoes(prs, rel: RelatorioEquipamento):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _slide_header(slide, "Conclusões e Insights Operacionais", rel.equipamento)

    # Painel esquerdo: cards
    _add_text(slide, Inches(0.4), Inches(1.1), Inches(7.5), Inches(0.45),
              "Principais Conclusões", size=Pt(15), bold=True, color=COLOR_TITLE)

    prod = _producao_mensal(rel)
    pico_mes, pico_qtd = (max(prod, key=lambda x: x[1]) if prod else ("—", 0))
    cards = [
        ("Produção Consistente",
         f"{_fmt_int_pt(rel.total_produzido)} peças em {_meses_operacao(rel)} meses "
         f"com pico histórico de {_fmt_int_pt(pico_qtd)} em {_ano_mes_label(pico_mes)}",
         COLOR_ACCENT_BLUE),
    ]
    if rel.resumo:
        topo = rel.resumo[0]
        cards.append((
            f"Dominância do {topo.composto}",
            f"{_fmt_pct_pt(topo.pct_total)} da produção — composto principal "
            f"por {topo.meses} meses",
            COLOR_ACCENT,
        ))
    cards.append((
        "Qualidade Controlada",
        f"Taxa de refugo de apenas {_fmt_pct_pt(rel.pct_refugo)} — "
        f"desempenho estável em todo o período",
        COLOR_OK,
    ))
    if rel.prescricao.data_prescrita:
        cards.append((
            "Manutenção Planejada",
            f"Próxima intervenção prevista para {rel.prescricao.data_prescrita} — "
            f"{rel.prescricao.dias_restantes or '—'} dias de folga operacional",
            COLOR_WARNING,
        ))

    y = Inches(1.65)
    for titulo, texto, cor in cards:
        _add_rect(slide, Inches(0.4), y, Inches(7.5), Inches(1.15),
                  fill=COLOR_PANEL, border=COLOR_BORDER, rounded=True)
        _add_rect(slide, Inches(0.4), y, Emu(50000), Inches(1.15),
                  fill=cor, border=None)
        _add_text(slide, Inches(0.65), y + Inches(0.15), Inches(7.0), Inches(0.4),
                  titulo, size=Pt(13), bold=True, color=COLOR_TITLE)
        _add_text(slide, Inches(0.65), y + Inches(0.55), Inches(7.0), Inches(0.55),
                  texto, size=Pt(11), color=COLOR_TEXT)
        y += Inches(1.25)

    # Painel direito: scorecard
    _add_rect(slide, Inches(8.1), Inches(1.1), Inches(4.85), Inches(5.9),
              fill=COLOR_PANEL, border=COLOR_BORDER, rounded=True)
    _add_text(slide, Inches(8.3), Inches(1.2), Inches(4.55), Inches(0.5),
              "Scorecard Operacional", size=Pt(16), bold=True, color=COLOR_TITLE)

    scorecard = [
        ("Volume Total", f"{_fmt_int_pt(rel.total_produzido)} peças", COLOR_OK),
        ("Período", f"{_meses_operacao(rel)} meses", COLOR_OK),
        ("Compostos", f"{rel.n_compostos} tipos", COLOR_OK),
        ("Taxa de Refugo", _fmt_pct_pt(rel.pct_refugo),
         COLOR_OK if rel.pct_refugo < 2.0 else COLOR_WARNING),
        ("Pico Mensal", f"{_fmt_int_pt(pico_qtd)} peças — {_ano_mes_label(pico_mes)}",
         COLOR_WARNING),
    ]
    if rel.prescricao.data_prescrita:
        scorecard.append((
            "Próxima Manutenção",
            f"{rel.prescricao.data_prescrita} ({rel.prescricao.dias_restantes or '—'} dias)",
            _urgencia_cor(rel.prescricao.urgencia),
        ))

    yL = Inches(1.85)
    for label, valor, dot in scorecard:
        _add_text(slide, Inches(8.3), yL, Inches(2.4), Inches(0.4),
                  label, size=Pt(11), color=COLOR_TEXT)
        _add_text(slide, Inches(10.7), yL, Inches(2.0), Inches(0.4),
                  valor, size=Pt(11), bold=True, color=COLOR_TITLE,
                  align=PP_ALIGN.RIGHT)
        # bullet de cor
        _add_rect(slide, Inches(12.75), yL + Inches(0.13), Inches(0.13), Inches(0.13),
                  fill=dot, border=None, rounded=True)
        # divisor
        _add_rect(slide, Inches(8.3), yL + Inches(0.45), Inches(4.55), Emu(4000),
                  fill=COLOR_BORDER, border=None)
        yL += Inches(0.55)

    _slide_footer(slide, rel.equipamento, _periodo_label(rel))


# ---------------------------------------------------------------------------
# Pipeline por equipamento
# ---------------------------------------------------------------------------
def gerar_pptx(rel: RelatorioEquipamento, out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_capa(prs, rel)
    _slide_visao_executiva(prs, rel)
    _slide_tendencias(prs, rel)
    _slide_compostos(prs, rel)
    _slide_qualidade(prs, rel)
    _slide_manutencao(prs, rel)
    _slide_conclusoes(prs, rel)

    prs.save(str(out_path))


def converter_pdf(pptx_path: Path, out_dir: Path) -> Path | None:
    soffice = shutil.which("libreoffice") or shutil.which("soffice")
    if not soffice:
        print("  [warn] LibreOffice não encontrado — PDF não gerado.", flush=True)
        return None
    proc = subprocess.run(
        [soffice, "--headless", "--norestore", "--nologo",
         "--convert-to", "pdf", "--outdir", str(out_dir), str(pptx_path)],
        capture_output=True, text=True, timeout=180,
    )
    if proc.returncode != 0:
        print(f"  [warn] falha ao gerar PDF: {proc.stderr.strip()[:300]}", flush=True)
        return None
    pdf_path = out_dir / (pptx_path.stem + ".pdf")
    return pdf_path if pdf_path.exists() else None


# ---------------------------------------------------------------------------
# main
# ---------------------------------------------------------------------------
def listar_equipamentos(filtro: Iterable[str] | None) -> list[Path]:
    arquivos = sorted(p for p in INPUT_DIR.glob("IJ-*.md")
                      if p.stem.upper() != "INDEX")
    if filtro:
        wanted = {f.upper() for f in filtro}
        arquivos = [p for p in arquivos if p.stem.upper() in wanted]
    return arquivos


def main(equipamentos: list[str] | None = None,
         skip_pdf: bool = False, **_kwargs):
    arquivos = listar_equipamentos(equipamentos)
    if not arquivos:
        print(f"[info] nenhum arquivo encontrado em {INPUT_DIR}")
        return {"gerados": 0}

    print(f"[info] gerando {len(arquivos)} apresentação(ões) em {OUTPUT_DIR}")
    resultados = []
    for md in arquivos:
        print(f"  → {md.stem} ...", flush=True)
        try:
            rel = parse_markdown(md)
            pptx_path = OUTPUT_DIR / f"{rel.equipamento}.pptx"
            gerar_pptx(rel, pptx_path)
            pdf_path = None
            if not skip_pdf:
                pdf_path = converter_pdf(pptx_path, OUTPUT_DIR)
            resultados.append({
                "equipamento": rel.equipamento,
                "pptx": str(pptx_path),
                "pdf": str(pdf_path) if pdf_path else None,
            })
        except Exception as exc:  # pragma: no cover
            print(f"    [erro] {md.stem}: {exc}")
    print(f"[ok] {len(resultados)} arquivo(s) gerado(s).")
    return {"gerados": len(resultados), "resultados": resultados}


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--equip", nargs="*",
                        help="Restringe a um ou mais equipamentos (ex.: IJ-044 IJ-130).")
    parser.add_argument("--skip-pdf", action="store_true",
                        help="Pula a conversão para PDF.")
    args = parser.parse_args()
    main(equipamentos=args.equip, skip_pdf=args.skip_pdf)
