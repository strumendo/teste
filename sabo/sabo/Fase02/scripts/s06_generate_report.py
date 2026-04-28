"""
S06 - Geração de Relatório PDF
==============================
Etapa 6 do Pipeline - Geração do Relatório de Entrega

O QUE FAZ:
- Gera relatório PDF completo no formato padrão SABO
- Inclui todas as seções: Resumo, Metodologia, Modelos, Métricas, Resultados
- Incorpora gráficos gerados nas etapas anteriores
- Documenta variáveis, desempenho e recomendações

ENTRADA:
- best_model.joblib (modelo selecionado)
- evaluation_report.txt (métricas)
- eda_report.txt (estatísticas)
- eda_plots/*.png (gráficos)
- data_eda.csv (dados processados)

SAÍDA:
- Relatorio_SABO_RX.pdf: Relatório completo de entrega
"""

import os
import json
from pathlib import Path
from datetime import datetime, timedelta
import warnings
warnings.filterwarnings('ignore')

# Verificar disponibilidade de bibliotecas para PDF
try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm, mm
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
        PageBreak, Image, ListFlowable, ListItem
    )
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False
    print("⚠ ReportLab não disponível. Gerando relatório em formato texto.")

import pandas as pd
import numpy as np
import joblib

try:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import matplotlib.ticker as mticker
    HAS_MATPLOTLIB = True
except ImportError:
    HAS_MATPLOTLIB = False
    print("⚠ Matplotlib não disponível. Gráficos mensais não serão gerados.")


# =============================================================================
# DICIONÁRIO COMPLETO DE VARIÁVEIS DO PIPELINE
# =============================================================================

VARIAVEIS_PIPELINE = {
    # =========================================================================
    # VARIÁVEIS ORIGINAIS (DADOS BRUTOS) - Etapa 1: s01_data_collection.py
    # =========================================================================
    "dados_brutos": {
        "titulo": "Variáveis Originais (Dados Brutos)",
        "descricao": "Variáveis coletadas diretamente dos arquivos CSV/XLSX dos equipamentos na Etapa 1.",
        "variaveis": [
            {
                "nome": "Data de Produção",
                "nome_processado": "Data_de_Producao",
                "tipo": "datetime",
                "descricao": "Data em que a produção foi realizada no equipamento.",
                "exemplo": "01/01/2024",
                "origem": "Arquivo CSV do equipamento",
                "uso": "Cálculo temporal, ordenação cronológica e cálculo da variável target."
            },
            {
                "nome": "Cód. Ordem",
                "nome_processado": "Cod_Ordem",
                "tipo": "int/string",
                "descricao": "Código identificador único do lote ou ordem de produção. Representa um conjunto específico de peças fabricadas em uma mesma operação.",
                "exemplo": "1234567",
                "origem": "Sistema de controle de produção (ERP)",
                "uso": "Rastreabilidade, agrupamento de lotes e análise de padrões por ordem."
            },
            {
                "nome": "Cód. Recurso",
                "nome_processado": "Cod_Recurso",
                "tipo": "string",
                "descricao": "Código do recurso produtivo (máquina/equipamento) utilizado na operação.",
                "exemplo": "IJ-044",
                "origem": "Sistema de controle de produção",
                "uso": "Identificação do equipamento para análise individualizada."
            },
            {
                "nome": "Cód. Produto",
                "nome_processado": "Cod_Produto",
                "tipo": "string",
                "descricao": "Código identificador único do produto fabricado. Cada código representa um tipo específico de peça de borracha.",
                "exemplo": "SA05780",
                "origem": "Cadastro de produtos",
                "uso": "Análise de consumo por produto, categorização para One-Hot Encoding."
            },
            {
                "nome": "Qtd. Produzida",
                "nome_processado": "Qtd_Produzida",
                "tipo": "int",
                "descricao": "Quantidade total de peças produzidas no registro. Representa o volume de produção efetiva do lote.",
                "exemplo": "1500",
                "unidade": "peças",
                "origem": "Apontamento de produção",
                "uso": "Cálculo de acumulados, análise de desgaste, correlação com manutenção."
            },
            {
                "nome": "Qtd. Refugada",
                "nome_processado": "Qtd_Refugada",
                "tipo": "int",
                "descricao": "Quantidade de peças refugadas (descartadas) por não atenderem aos padrões de qualidade. Indicador de qualidade do processo.",
                "exemplo": "25",
                "unidade": "peças",
                "origem": "Controle de qualidade",
                "uso": "Indicador de qualidade, correlação com necessidade de manutenção."
            },
            {
                "nome": "Qtd. Retrabalhada",
                "nome_processado": "Qtd_Retrabalhada",
                "tipo": "int",
                "descricao": "Quantidade de peças que necessitaram de retrabalho antes de serem aprovadas. Indica problemas intermediários no processo.",
                "exemplo": "10",
                "unidade": "peças",
                "origem": "Controle de qualidade",
                "uso": "Indicador de eficiência, correlação com condição do equipamento."
            },
            {
                "nome": "Fator Un.",
                "nome_processado": "Fator_Un",
                "tipo": "float",
                "descricao": "Fator de conversão de unidade. Usado para normalização de quantidades quando unidades diferem.",
                "exemplo": "1.0",
                "origem": "Cadastro de produtos",
                "uso": "Normalização de quantidades entre diferentes unidades de medida."
            },
            {
                "nome": "Cód. Un.",
                "nome_processado": "Cod_Un",
                "tipo": "string",
                "descricao": "Código da unidade de medida utilizada para quantificar a produção.",
                "exemplo": "PC (peças)",
                "origem": "Cadastro de produtos",
                "uso": "Identificação da unidade, categorização."
            },
            {
                "nome": "Descrição da massa (Composto)",
                "nome_processado": "Descricao_da_massa_Composto",
                "tipo": "string",
                "descricao": "Tipo de composto ou material de borracha utilizado na fabricação. Diferentes compostos têm diferentes propriedades de desgaste.",
                "exemplo": "N-142/67",
                "origem": "Especificação técnica do produto",
                "uso": "Análise de desgaste por tipo de material, One-Hot Encoding."
            },
            {
                "nome": "Consumo de massa no item em (Kg/100pçs)",
                "nome_processado": "Consumo_de_massa_no_item_em_Kg_100pcs",
                "tipo": "float",
                "descricao": "Consumo de matéria-prima (borracha) por cada 100 peças produzidas. Indicador de eficiência de material.",
                "exemplo": "1.250",
                "unidade": "Kg/100 peças",
                "origem": "Cálculo baseado em consumo real",
                "uso": "Análise de eficiência, detecção de anomalias de consumo."
            },
            {
                "nome": "Equipamento",
                "nome_processado": "Equipamento",
                "tipo": "string",
                "descricao": "Identificador do equipamento/extrusora. Extraído do nome do arquivo de origem (IJ-XXX).",
                "exemplo": "IJ-044",
                "origem": "Nome do arquivo CSV",
                "uso": "Agrupamento por equipamento, One-Hot Encoding, cálculo de troca de peças."
            },
            {
                "nome": "Fonte_Dados",
                "nome_processado": "Fonte_Dados",
                "tipo": "string",
                "descricao": "Nome do arquivo de origem dos dados. Usado para rastreabilidade em arquivos com formato estendido.",
                "exemplo": "IJ-138.2.xlsx",
                "origem": "Sistema de arquivos",
                "uso": "Rastreabilidade e auditoria de dados."
            },
        ]
    },

    # =========================================================================
    # VARIÁVEIS GERADAS NO PRÉ-PROCESSAMENTO - Etapa 2: s02_preprocessing.py
    # =========================================================================
    "preprocessamento": {
        "titulo": "Variáveis Geradas no Pré-processamento",
        "descricao": "Variáveis criadas durante a Etapa 2 através de transformações e engenharia de features.",
        "variaveis": [
            {
                "nome": "Manutencao",
                "nome_processado": "Manutencao",
                "tipo": "int",
                "descricao": "VARIÁVEL TARGET - Número de dias restantes até a próxima manutenção programada do equipamento. Calculada como a diferença entre a data de manutenção conhecida e a data de produção.",
                "exemplo": "45",
                "unidade": "dias",
                "origem": "Calculada (data_manutencao - data_producao)",
                "uso": "Variável alvo para os modelos de Machine Learning.",
                "formula": "Manutencao = Data_Manutenção_Equipamento - Data_de_Produção"
            },
            {
                "nome": "Qtd. Produzida_Acumulado",
                "nome_processado": "Qtd_Produzida_Acumulado",
                "tipo": "float",
                "descricao": "Quantidade acumulada de peças produzidas por equipamento desde o início do período. Soma cumulativa ordenada por data.",
                "exemplo": "150000",
                "unidade": "peças",
                "origem": "Soma cumulativa por equipamento",
                "uso": "Indicador de desgaste acumulado, correlação com necessidade de manutenção.",
                "formula": "cumsum(Qtd_Produzida) agrupado por Equipamento"
            },
            {
                "nome": "Qtd. Refugada_Acumulado",
                "nome_processado": "Qtd_Refugada_Acumulado",
                "tipo": "float",
                "descricao": "Quantidade acumulada de peças refugadas por equipamento. Indica a evolução da taxa de refugo ao longo do tempo.",
                "exemplo": "2500",
                "unidade": "peças",
                "origem": "Soma cumulativa por equipamento",
                "uso": "Tendência de qualidade, detecção de degradação.",
                "formula": "cumsum(Qtd_Refugada) agrupado por Equipamento"
            },
            {
                "nome": "Qtd. Retrabalhada_Acumulado",
                "nome_processado": "Qtd_Retrabalhada_Acumulado",
                "tipo": "float",
                "descricao": "Quantidade acumulada de peças retrabalhadas por equipamento. Indica a evolução de problemas de processo.",
                "exemplo": "1200",
                "unidade": "peças",
                "origem": "Soma cumulativa por equipamento",
                "uso": "Tendência de eficiência, correlação com condição do equipamento.",
                "formula": "cumsum(Qtd_Retrabalhada) agrupado por Equipamento"
            },
            {
                "nome": "Consumo de massa_Acumulado",
                "nome_processado": "Consumo_de_massa_no_item_em_Kg_100pcs_Acumulado",
                "tipo": "float",
                "descricao": "Consumo acumulado de matéria-prima por equipamento. Permite análise de eficiência de material ao longo do tempo.",
                "exemplo": "5000.50",
                "unidade": "Kg",
                "origem": "Soma cumulativa por equipamento",
                "uso": "Análise de consumo, detecção de anomalias de material.",
                "formula": "cumsum(Consumo_de_massa) agrupado por Equipamento"
            },
        ]
    },

    # =========================================================================
    # VARIÁVEIS DE ENCODING (ONE-HOT) - Etapa 2: s02_preprocessing.py
    # =========================================================================
    "encoding": {
        "titulo": "Variáveis de Encoding (One-Hot)",
        "descricao": "Variáveis binárias criadas através de One-Hot Encoding para representar categorias. Cada categoria única gera uma coluna binária (0 ou 1).",
        "variaveis": [
            {
                "nome": "Equipamento_*",
                "nome_processado": "Equipamento_IJ-XXX",
                "tipo": "int (0/1)",
                "descricao": "Variáveis binárias para cada equipamento. Valor 1 indica que o registro pertence àquele equipamento específico.",
                "exemplo": "Equipamento_IJ-044 = 1",
                "origem": "One-Hot Encoding de 'Equipamento'",
                "uso": "Permitir que modelos capturem comportamentos específicos de cada equipamento.",
                "categorias": "IJ-044, IJ-046, IJ-117, IJ-118, IJ-119, IJ-120, IJ-121, IJ-122, IJ-123, IJ-124, IJ-125, IJ-129, IJ-130, IJ-131, IJ-132, IJ-133, IJ-134, IJ-135, IJ-136, IJ-137, IJ-138, IJ-139, IJ-151, IJ-152, IJ-155, IJ-156, IJ-164"
            },
            {
                "nome": "Cod_Produto_*",
                "nome_processado": "Cod_Produto_SAXXXXX",
                "tipo": "int (0/1)",
                "descricao": "Variáveis binárias para cada código de produto. Permite análise de impacto por tipo de produto.",
                "exemplo": "Cod_Produto_SA05780 = 1",
                "origem": "One-Hot Encoding de 'Cod_Produto'",
                "uso": "Capturar diferenças de desgaste entre produtos.",
                "categorias": "Variável conforme produtos únicos no dataset"
            },
            {
                "nome": "Descricao_da_massa_*",
                "nome_processado": "Descricao_da_massa_Composto_*",
                "tipo": "int (0/1)",
                "descricao": "Variáveis binárias para cada tipo de composto/material. Diferentes compostos causam diferentes níveis de desgaste.",
                "exemplo": "Descricao_da_massa_Composto_N-142/67 = 1",
                "origem": "One-Hot Encoding de 'Descrição da massa'",
                "uso": "Análise de impacto de diferentes materiais no desgaste.",
                "categorias": "Variável conforme compostos únicos no dataset"
            },
            {
                "nome": "Cod_Un_*",
                "nome_processado": "Cod_Un_*",
                "tipo": "int (0/1)",
                "descricao": "Variáveis binárias para cada unidade de medida.",
                "exemplo": "Cod_Un_PC = 1",
                "origem": "One-Hot Encoding de 'Cód. Un.'",
                "uso": "Normalização entre diferentes unidades.",
                "categorias": "PC, UN, etc."
            },
        ]
    },

    # =========================================================================
    # HISTÓRICO DE TROCA DE PEÇAS POR EQUIPAMENTO
    # (Carregado dinamicamente de equipment_stats.json)
    # =========================================================================
    "troca_pecas_equipamentos": {
        "titulo": "Histórico de Troca de Peças por Equipamento",
        "descricao": "Datas das trocas efetivas de peças (substituição de componentes) e intervalo de operação para cada equipamento.",
        "equipamentos": {}  # Será preenchido dinamicamente
    }
}


def predict_maintenance_with_ml(equip_data: dict) -> dict:
    """
    Usa o modelo ML treinado para prever dias até manutenção para cada equipamento.

    O modelo foi treinado com features como produção acumulada, desgaste, refugo, etc.
    Para cada equipamento, usa o último estado conhecido para fazer a previsão.

    Args:
        equip_data: Dicionário com dados de equipamentos do equipment_stats.json

    Returns:
        Dicionário {equipamento: dias_previstos_ml}
    """
    predictions = {}

    # Carregar modelo (preferir o mais recente)
    model_path = Path("best_model.joblib")
    if not model_path.exists():
        model_path = Path("models/best_model.joblib")

    if not model_path.exists():
        print("  ⚠ Modelo ML não encontrado. Previsões ML indisponíveis.")
        return predictions

    # Carregar dados de treino para obter nomes das features
    train_data_path = Path("train_test_split.npz")
    if not train_data_path.exists():
        print("  ⚠ Dados de treino não encontrados. Previsões ML indisponíveis.")
        return predictions

    try:
        model_data = joblib.load(model_path)
        # O modelo pode estar salvo diretamente ou dentro de um dicionário
        if isinstance(model_data, dict) and 'model' in model_data:
            model = model_data['model']
        else:
            model = model_data

        train_data = np.load(train_data_path, allow_pickle=True)
        feature_names = list(train_data['feature_names'])

        # Carregar dados EDA para obter último registro de cada equipamento
        eda_path = Path("data_eda.csv")
        if not eda_path.exists():
            print("  ⚠ data_eda.csv não encontrado. Previsões ML indisponíveis.")
            return predictions

        df_eda = pd.read_csv(eda_path)

        # Para cada equipamento, pegar o último registro e fazer previsão
        for equip in equip_data.keys():
            # Encontrar coluna do equipamento (One-Hot encoded)
            equip_col = f"Equipamento_{equip.replace('-', '_')}"

            if equip_col in df_eda.columns:
                # Filtrar registros deste equipamento
                df_equip = df_eda[df_eda[equip_col] == True]

                if len(df_equip) > 0:
                    # Pegar último registro (estado mais recente)
                    last_row = df_equip.iloc[-1]

                    # Preparar features para previsão
                    X_pred = []
                    for feat in feature_names:
                        if feat in last_row.index:
                            val = last_row[feat]
                            # Tratar NaN
                            if pd.isna(val):
                                val = 0.0
                            X_pred.append(float(val))
                        else:
                            X_pred.append(0.0)

                    # Fazer previsão
                    X_pred = np.array([X_pred])
                    pred = model.predict(X_pred)[0]

                    # Garantir valor positivo e arredondar
                    predictions[equip] = max(0, int(round(pred)))

        print(f"  ✓ Previsões ML geradas para {len(predictions)} equipamentos")

    except Exception as e:
        print(f"  ⚠ Erro ao gerar previsões ML: {e}")

    return predictions


def _load_sap_scheduled_dates() -> dict:
    """
    Lê o XLSX 'Histórico Geral Preventivas RM.195' e devolve, por equipamento,
    a lista de datas agendadas (ordem ascendente) como pd.Timestamp.
    """
    candidates = [
        Path("../data/manutencao"),
        Path("data/manutencao"),
        Path(__file__).resolve().parent.parent / "data" / "manutencao",
    ]
    fp = None
    for d in candidates:
        if d.exists():
            matches = list(d.glob("Histórico Geral Preventivas*.xlsx"))
            if matches:
                fp = matches[0]
                break
    if fp is None:
        print("  ⚠ Arquivo de preventivas RM.195 não encontrado para SAP")
        return {}
    try:
        df = pd.read_excel(fp)
    except Exception as e:
        print(f"  ⚠ Erro ao ler preventivas: {e}")
        return {}
    equip_col = next((c for c in df.columns if c.strip().lower() == "equipamento"), None)
    date_col = next((c for c in df.columns if "iníc" in c.lower() or "inic" in c.lower()), None)
    if equip_col is None or date_col is None:
        return {}
    df = df[[equip_col, date_col]].copy()
    df[date_col] = pd.to_datetime(df[date_col], errors="coerce")
    df = df.dropna(subset=[date_col])
    out: dict = {}
    for equip, grp in df.groupby(equip_col):
        out[str(equip).strip()] = sorted(grp[date_col].tolist())
    return out


def _load_ml_predictions() -> dict:
    """
    Carrega outputs/prescricao_manutencao.csv e retorna {equip: pd.Timestamp data_prescrita}.
    """
    fp = Path("prescricao_manutencao.csv")
    if not fp.exists():
        return {}
    try:
        df = pd.read_csv(fp)
    except Exception as e:
        print(f"  ⚠ Erro ao ler prescricao_manutencao.csv: {e}")
        return {}
    if "equipamento" not in df.columns or "data_prescrita" not in df.columns:
        return {}
    df["data_prescrita"] = pd.to_datetime(df["data_prescrita"], errors="coerce", dayfirst=True)
    out = {}
    for _, r in df.iterrows():
        if pd.notna(r["data_prescrita"]):
            out[str(r["equipamento"]).strip()] = r["data_prescrita"]
    return out


def _load_last_production_dates() -> dict:
    """
    Lê outputs/data_raw.csv e devolve {equip: pd.Timestamp última_data_producao}.
    """
    fp = Path("data_raw.csv")
    if not fp.exists():
        return {}
    try:
        df = pd.read_csv(fp, usecols=["Data de Produção", "Equipamento"])
    except Exception as e:
        print(f"  ⚠ Erro ao ler data_raw.csv: {e}")
        return {}
    df["Data de Produção"] = pd.to_datetime(df["Data de Produção"], errors="coerce")
    df = df.dropna(subset=["Data de Produção", "Equipamento"])
    return df.groupby("Equipamento")["Data de Produção"].max().to_dict()


def generate_previsao_manutencao_csv(equip_data: dict, output_path: str = "previsao_manutencao_consolidada.csv") -> str:
    """
    Gera CSV consolidado de previsão de manutenção por equipamento.

    Colunas:
      Equipamento, Ultima_Data_Producao, Ultima_Manutencao,
      Intervalo_Penultima_Ultima_Manutencao_dias,
      Data_Proxima_Manutencao_SAP, Dias_Ultima_Manutencao_ate_SAP,
      Dias_SAP_ate_Hoje, Data_Proxima_Manutencao_ML,
      Diferenca_ML_vs_SAP_dias
    """
    today = pd.Timestamp(datetime.now().date())
    sap = _load_sap_scheduled_dates()
    ml = _load_ml_predictions()
    last_prod = _load_last_production_dates()

    rows = []
    for equip, dados in sorted(equip_data.items()):
        if not isinstance(dados, dict):
            continue
        ultima_str = dados.get("data_ultima_manutencao") or dados.get("ultima_troca")
        penultima_str = dados.get("data_penultima_manutencao")
        try:
            ultima_dt = pd.to_datetime(ultima_str) if ultima_str else pd.NaT
        except Exception:
            ultima_dt = pd.NaT
        try:
            penultima_dt = pd.to_datetime(penultima_str) if penultima_str else pd.NaT
        except Exception:
            penultima_dt = pd.NaT
        intervalo_dias = (ultima_dt - penultima_dt).days if pd.notna(ultima_dt) and pd.notna(penultima_dt) else None

        prod_dt = last_prod.get(equip)
        if prod_dt is None:
            for k, v in last_prod.items():
                if str(k).strip() == equip:
                    prod_dt = v
                    break

        sap_list = sap.get(equip, [])
        sap_proxima = next((d for d in sap_list if d >= today), None)

        ml_dt = ml.get(equip)

        dias_ult_ate_sap = (sap_proxima - ultima_dt).days if (sap_proxima is not None and pd.notna(ultima_dt)) else None
        dias_sap_hoje = (sap_proxima - today).days if sap_proxima is not None else None
        dif_ml_sap = (ml_dt - sap_proxima).days if (sap_proxima is not None and ml_dt is not None) else None

        rows.append({
            "Equipamento": equip,
            "Ultima_Data_Producao": prod_dt.strftime("%Y-%m-%d") if isinstance(prod_dt, pd.Timestamp) and pd.notna(prod_dt) else "",
            "Ultima_Manutencao": ultima_dt.strftime("%Y-%m-%d") if pd.notna(ultima_dt) else "",
            "Intervalo_Penultima_Ultima_Manutencao_dias": intervalo_dias if intervalo_dias is not None else "",
            "Data_Proxima_Manutencao_SAP": sap_proxima.strftime("%Y-%m-%d") if sap_proxima is not None else "",
            "Dias_Ultima_Manutencao_ate_SAP": dias_ult_ate_sap if dias_ult_ate_sap is not None else "",
            "Dias_SAP_ate_Hoje": dias_sap_hoje if dias_sap_hoje is not None else "",
            "Data_Proxima_Manutencao_ML": ml_dt.strftime("%Y-%m-%d") if ml_dt is not None else "",
            "Diferenca_ML_vs_SAP_dias": dif_ml_sap if dif_ml_sap is not None else "",
        })

    df = pd.DataFrame(rows)
    df.to_csv(output_path, index=False, encoding="utf-8")
    print(f"  ✓ CSV consolidado salvo: {output_path} ({len(df)} equipamentos)")
    return output_path


def merge_componentes_pptx(output_path: str = "Apresentacoes_Componentes_Consolidado.pptx") -> str | None:
    """
    Concatena todos os PPTX 'IJ-*.pptx' em outputs/relatorios_mensais_componentes_ppt/
    em um único arquivo, usando python-pptx para clonar slides com seus relacionamentos.
    """
    try:
        from pptx import Presentation
        from pptx.util import Emu
        from copy import deepcopy
    except ImportError:
        print("  ⚠ python-pptx não disponível. Pulando consolidação de PPTX.")
        return None

    src_dir = Path("relatorios_mensais_componentes_ppt")
    if not src_dir.exists():
        print(f"  ⚠ Diretório {src_dir} não encontrado")
        return None

    pptx_files = sorted(src_dir.glob("IJ-*.pptx"))
    if not pptx_files:
        print(f"  ⚠ Nenhum PPTX 'IJ-*' em {src_dir}")
        return None

    base = Presentation(str(pptx_files[0]))
    base_slide_w = base.slide_width
    base_slide_h = base.slide_height

    def _copy_slide(target_pres, src_slide):
        blank_layout = target_pres.slide_layouts[6] if len(target_pres.slide_layouts) > 6 else target_pres.slide_layouts[-1]
        new_slide = target_pres.slides.add_slide(blank_layout)
        for shape in new_slide.shapes:
            pass
        for shp in list(new_slide.shapes):
            sp = shp._element
            sp.getparent().remove(sp)
        for shp in src_slide.shapes:
            el = shp.element
            new_slide.shapes._spTree.insert_element_before(deepcopy(el), 'p:extLst')
        for rel in src_slide.part.rels.values():
            if "image" in rel.reltype or "chart" in rel.reltype or "media" in rel.reltype:
                try:
                    new_slide.part.relate_to(rel.target_part, rel.reltype)
                except Exception:
                    pass
        return new_slide

    for src_path in pptx_files[1:]:
        try:
            src = Presentation(str(src_path))
            for slide in src.slides:
                _copy_slide(base, slide)
        except Exception as e:
            print(f"  ⚠ Erro ao mesclar {src_path.name}: {e}")
            continue

    base.slide_width = base_slide_w
    base.slide_height = base_slide_h
    out = Path(output_path)
    base.save(str(out))
    print(f"  ✓ PPTX consolidado salvo: {out} ({len(pptx_files)} arquivos mesclados)")
    return str(out)


def load_equipment_stats() -> dict:
    """
    Carrega estatísticas de equipamentos do arquivo equipment_stats.json.

    Returns:
        Dicionário com dados de cada equipamento no formato:
        {"IJ-XXX": {"ultima_troca": "YYYY-MM-DD", "intervalo_dias": N, ...}}
    """
    stats_file = Path("equipment_stats.json")

    if not stats_file.exists():
        print("  ⚠ equipment_stats.json não encontrado. Usando dados vazios.")
        return {}

    try:
        with open(stats_file, "r", encoding="utf-8") as f:
            stats_list = json.load(f)

        # Converter lista para dicionário indexado por equipamento
        equip_data = {}
        for item in stats_list:
            equip = item.get("equipamento")
            if equip:
                equip_data[equip] = {
                    "ultima_troca": item.get("data_ultima_manutencao"),
                    "data_ultima_manutencao": item.get("data_ultima_manutencao"),
                    "data_penultima_manutencao": item.get("data_penultima_manutencao"),
                    "penultima_troca": item.get("data_penultima_manutencao"),
                    "intervalo_dias": item.get("intervalo_manutencao_dias", 365),
                    "intervalo_manutencao_dias": item.get("intervalo_manutencao_dias"),
                    "media_dias_manutencao": item.get("media_dias_manutencao"),
                    "min_dias_manutencao": item.get("min_dias_manutencao"),
                    "max_dias_manutencao": item.get("max_dias_manutencao"),
                    "total_produzido": item.get("total_produzido", 0),
                    "total_refugado": item.get("total_refugado", 0),
                    "taxa_refugo_pct": item.get("taxa_refugo_pct", 0),
                    "indice_desgaste": item.get("indice_desgaste_medio", 0),
                    "cilindro_max": item.get("cilindro_max"),
                    "cilindro_min": item.get("cilindro_min"),
                    "fuso_max": item.get("fuso_max"),
                    "fuso_min": item.get("fuso_min"),
                    "observacoes": item.get("observacoes_manutencao"),
                }

        print(f"  ✓ Carregados dados de {len(equip_data)} equipamentos")
        return equip_data

    except Exception as e:
        print(f"  ⚠ Erro ao carregar equipment_stats.json: {e}")
        return {}


def get_report_version(suffix: str = "") -> str:
    """
    Determina a versão do relatório baseado em arquivos existentes.

    Args:
        suffix: Sufixo opcional para o nome (ex: "_v1")

    Returns:
        String com versão (R5, R6, etc.) + sufixo
    """
    existing = list(Path(".").glob("Relatorio_SABO_R*.pdf"))
    if not existing:
        return f"R5{suffix}"

    versions = []
    for f in existing:
        try:
            # Extrair número da versão, ignorando sufixos
            version_part = f.stem.split("_R")[-1]
            # Remover sufixos não numéricos (ex: "_v1")
            num_str = ""
            for c in version_part:
                if c.isdigit():
                    num_str += c
                else:
                    break
            if num_str:
                versions.append(int(num_str))
        except:
            pass

    if versions:
        return f"R{max(versions) + 1}{suffix}"
    return f"R5{suffix}"


def load_pipeline_results() -> dict:
    """
    Carrega todos os resultados das etapas anteriores do pipeline.

    Returns:
        Dicionário com todos os dados para o relatório
    """
    results = {
        "data": {},
        "model": {},
        "metrics": {},
        "eda": {},
        "plots": [],
    }

    # Carregar modelo
    if Path("best_model.joblib").exists():
        model_data = joblib.load("best_model.joblib")
        if isinstance(model_data, dict):
            results["model"] = model_data
        else:
            results["model"] = {"model": model_data, "name": "unknown"}

    # Carregar dados
    if Path("data_eda.csv").exists():
        df = pd.read_csv("data_eda.csv")
        results["data"]["shape"] = df.shape
        results["data"]["columns"] = list(df.columns)
        results["data"]["dtypes"] = df.dtypes.to_dict()

        # Estatísticas básicas
        numeric_cols = df.select_dtypes(include=[np.number]).columns
        results["data"]["stats"] = df[numeric_cols].describe().to_dict()

        # Extrair período dos dados (procurar colunas de data)
        date_columns = [col for col in df.columns if 'data' in col.lower() or 'date' in col.lower() or 'dt' in col.lower()]
        if date_columns:
            for date_col in date_columns:
                try:
                    df[date_col] = pd.to_datetime(df[date_col], errors='coerce')
                    valid_dates = df[date_col].dropna()
                    if len(valid_dates) > 0:
                        results["data"]["periodo_inicio"] = valid_dates.min().strftime('%d/%m/%Y')
                        results["data"]["periodo_fim"] = valid_dates.max().strftime('%d/%m/%Y')
                        results["data"]["total_dias"] = (valid_dates.max() - valid_dates.min()).days
                        break
                except:
                    pass

        # Calcular total de peças produzidas
        prod_columns = [col for col in df.columns if 'produzid' in col.lower() or 'quantidade' in col.lower() or 'qtd' in col.lower()]
        if prod_columns:
            for prod_col in prod_columns:
                try:
                    results["data"]["total_pecas"] = int(df[prod_col].sum())
                    results["data"]["coluna_producao"] = prod_col
                    break
                except:
                    pass

    # Carregar relatório de EDA
    if Path("eda_report.txt").exists():
        with open("eda_report.txt", "r", encoding="utf-8") as f:
            results["eda"]["report"] = f.read()

    # Carregar relatório de avaliação
    if Path("evaluation_report.txt").exists():
        with open("evaluation_report.txt", "r", encoding="utf-8") as f:
            results["metrics"]["report"] = f.read()

    # Listar gráficos
    plots_dir = Path("eda_plots")
    if plots_dir.exists():
        results["plots"] = list(plots_dir.glob("*.png"))

    # Carregar dados de treino/teste
    if Path("train_test_split.npz").exists():
        data = np.load("train_test_split.npz", allow_pickle=True)
        results["data"]["train_size"] = len(data["X_train"])
        results["data"]["test_size"] = len(data["X_test"])
        results["data"]["features"] = data["feature_names"].tolist()

    # Carregar histórico de execuções para seção comparativa
    results["history"] = []
    history_index = Path("history/index.json")
    if history_index.exists():
        try:
            with open(history_index, "r", encoding="utf-8") as f:
                index_data = json.load(f)
            runs = index_data.get("runs", [])
            # Pegar últimas 10 execuções com dados válidos
            valid_runs = [r for r in runs if r.get("overall_best") is not None]
            for run_info in valid_runs[-10:]:
                run_id = run_info.get("run_id", "")
                run_file = Path(f"history/runs/run_{run_id}.json")
                if run_file.exists():
                    try:
                        with open(run_file, "r", encoding="utf-8") as f:
                            run_data = json.load(f)
                        # Extrair métricas do s05_evaluation
                        eval_step = run_data.get("steps", {}).get("s05_evaluation", {})
                        eval_models = eval_step.get("models", {})
                        best_model = eval_step.get("best_model", "")
                        best_r2 = eval_step.get("best_r2", None)
                        best_mse = eval_step.get("best_mse", None)

                        if best_r2 is not None and best_r2 != float('-inf'):
                            results["history"].append({
                                "run_id": run_id,
                                "timestamp": run_info.get("timestamp", ""),
                                "model": best_model,
                                "r2": best_r2,
                                "mse": best_mse,
                            })
                    except Exception:
                        pass
        except Exception as e:
            print(f"  Aviso: Erro ao carregar histórico - {e}")

    return results


def generate_pdf_report(results: dict, output_path: str, inicio: str = None, fim: str = None) -> str:
    """
    Gera relatório em PDF usando ReportLab.

    Args:
        results: Dados do pipeline
        output_path: Caminho do arquivo PDF
        inicio: Data de início do filtro de período (YYYY-MM-DD)
        fim: Data de fim do filtro de período (YYYY-MM-DD)

    Returns:
        Caminho do arquivo gerado
    """
    if not HAS_REPORTLAB:
        return generate_text_report(results, output_path.replace(".pdf", ".txt"))

    doc = SimpleDocTemplate(
        output_path,
        pagesize=A4,
        rightMargin=2*cm,
        leftMargin=2*cm,
        topMargin=2*cm,
        bottomMargin=2*cm
    )

    # Estilos
    styles = getSampleStyleSheet()

    # Estilos personalizados
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        spaceAfter=30,
        alignment=TA_CENTER,
        textColor=colors.HexColor('#333333')
    )

    heading1_style = ParagraphStyle(
        'CustomH1',
        parent=styles['Heading1'],
        fontSize=16,
        spaceBefore=20,
        spaceAfter=10,
        textColor=colors.HexColor('#F77F00')  # Laranja do logo Taking
    )

    heading2_style = ParagraphStyle(
        'CustomH2',
        parent=styles['Heading2'],
        fontSize=14,
        spaceBefore=15,
        spaceAfter=8,
        textColor=colors.HexColor('#333333')
    )

    body_style = ParagraphStyle(
        'CustomBody',
        parent=styles['Normal'],
        fontSize=11,
        alignment=TA_JUSTIFY,
        spaceAfter=12,
        leading=14
    )

    # Construir documento
    story = []

    # === CAPA ===
    story.append(Spacer(1, 3*cm))
    story.append(Paragraph("Relatório de Entrega", title_style))
    story.append(Spacer(1, 1*cm))
    story.append(Paragraph("SABO - Pipeline de Machine Learning", styles['Heading2']))
    story.append(Paragraph("Manutenção Prescritiva para Extrusoras de Borracha Y125", styles['Normal']))
    story.append(Spacer(1, 2*cm))
    story.append(Paragraph(f"Data: {datetime.now().strftime('%d/%m/%Y')}", styles['Normal']))
    story.append(PageBreak())

    # === SUMÁRIO ===
    story.append(Paragraph("Sumário", heading1_style))
    sumario_items = [
        "Resumo",
        "1. Período dos Dados Analisados",
        "2. Introdução",
        "3. Descrição do Método Utilizado",
        "4. Ferramentas Utilizadas",
        "5. Modelos Estocásticos e Estatísticos Utilizados",
        "6. Métricas",
        "7. Dicionário Completo de Variáveis do Pipeline",
        "   7.1 Variáveis Originais (Dados Brutos)",
        "   7.2 Variáveis Geradas no Pré-processamento",
        "   7.3 Variáveis de Encoding (One-Hot)",
        "   7.4 Histórico de Troca de Peças por Equipamento",
        "   7.5 Features Utilizadas no Modelo Final",
        "8. Dados Removidos da Análise",
        "9. Desempenho de Modelos",
        "10. Perguntas e Respostas",
        "11. Previsão de Troca de Peças",
        "12. Análise Mensal por Equipamento",
        "13. Considerações",
        "14. Próximos Passos",
        "15. Recomendações",
        "16. Considerações Finais",
        "17. Consumo de Massa por Equipamento e Composto",
        "Referência Bibliográfica"
    ]
    for item in sumario_items:
        story.append(Paragraph(f"• {item}", styles['Normal']))
    story.append(PageBreak())

    # === RESUMO ===
    story.append(Paragraph("Resumo", heading1_style))

    # Extrair métricas do modelo
    model_name = results.get("model", {}).get("name", "Não identificado")
    metrics = results.get("model", {}).get("metrics", {})
    r2 = metrics.get("r2", 0)
    mse = metrics.get("mse", 0)
    mae = metrics.get("mae", 0)

    data_shape = results.get("data", {}).get("shape", (0, 0))
    train_size = results.get("data", {}).get("train_size", 0)
    test_size = results.get("data", {}).get("test_size", 0)

    resumo_text = f"""
    Este relatório apresenta uma visão abrangente do projeto desenvolvido para prever a
    necessidade de manutenção em extrusoras de borracha modelo Y125, utilizando dados
    coletados em ambiente industrial. A iniciativa surgiu da observação de que, durante as
    manutenções preventivas, diversos componentes demonstraram uma vida útil
    significativamente maior que a inicialmente prevista pelo planejamento tradicional.

    Para viabilizar essa análise, foram coletados e organizados dados provenientes de
    múltiplos arquivos CSV, cada um associado a uma extrusora específica, contendo
    informações como quantidade produzida, refugada e retrabalhada, consumo de
    matéria-prima, datas de produção, códigos dos produtos e outros atributos operacionais.

    O dataset consolidado contém {data_shape[0]} registros e {data_shape[1]} variáveis.
    Foi realizada a divisão dos dados em conjuntos de treinamento ({train_size} amostras, 80%)
    e teste ({test_size} amostras, 20%) para avaliação do poder preditivo dos modelos.

    Os resultados indicaram desempenho promissor. O modelo {model_name.upper()} foi
    selecionado como melhor modelo, apresentando R² de {r2:.4f}, MSE de {mse:.2f} e
    MAE de {mae:.2f}.
    """
    story.append(Paragraph(resumo_text, body_style))
    story.append(PageBreak())

    # === 1. PERÍODO DOS DADOS ANALISADOS ===
    story.append(Paragraph("1. Período dos Dados Analisados", heading1_style))

    periodo_inicio = results.get("data", {}).get("periodo_inicio", "Não identificado")
    periodo_fim = results.get("data", {}).get("periodo_fim", "Não identificado")
    total_dias = results.get("data", {}).get("total_dias", 0)
    total_pecas = results.get("data", {}).get("total_pecas", 0)

    periodo_text = f"""
    <b>Período de Análise:</b> {periodo_inicio} a {periodo_fim}

    <b>Duração Total:</b> {total_dias} dias ({total_dias // 30} meses aproximadamente)

    <b>Total de Registros:</b> {data_shape[0]} registros de produção

    <b>Total de Peças Produzidas:</b> {total_pecas:,} peças
    """
    story.append(Paragraph(periodo_text, body_style))

    story.append(Paragraph("Base de Cálculo:", heading2_style))
    base_calculo_text = """
    Os cálculos e análises apresentados neste relatório utilizam as seguintes bases:

    • <b>Análise Temporal:</b> Baseada em dias corridos desde o início da coleta de dados
    • <b>Análise de Produção:</b> Baseada na quantidade total de peças produzidas por período
    • <b>Análise de Desgaste:</b> Correlação entre volume de produção e indicadores de manutenção
    • <b>Previsões:</b> Calculadas em função do número de peças produzidas e/ou dias de operação

    Os gráficos apresentados podem utilizar diferentes unidades de medida conforme indicado em suas
    respectivas legendas: dias, meses ou quantidade de peças.
    """
    story.append(Paragraph(base_calculo_text, body_style))

    # Informação de filtro de período (se aplicável)
    if inicio or fim:
        story.append(Paragraph("Filtro de Período Aplicado:", heading2_style))
        filtro_text = f"""
        Foi aplicado um filtro de período nesta execução do pipeline:
        <b>Data Início:</b> {inicio or 'Não definida (desde o início dos dados)'}
        <b>Data Fim:</b> {fim or 'Não definida (até o final dos dados)'}

        Apenas os registros dentro deste período foram considerados para análise,
        treinamento de modelos e geração de métricas.
        """
        story.append(Paragraph(filtro_text, body_style))

    story.append(PageBreak())

    # === 2. INTRODUÇÃO ===
    story.append(Paragraph("2. Introdução", heading1_style))
    intro_text = """
    A indústria moderna enfrenta o desafio constante de equilibrar segurança operacional e
    eficiência econômica. No contexto das extrusoras de borracha modelo Y125, observou-se
    que muitos componentes estavam sendo substituídos seguindo cronogramas fixos de
    manutenção preventiva, mesmo quando ainda apresentavam condições adequadas de
    funcionamento. Este cenário resulta em subutilização de peças, gerando custos
    desnecessários com componentes, mão de obra e tempo de máquina parada para
    intervenções prematuras.

    O objetivo deste projeto foi desenvolver e avaliar um modelo prescritivo capaz não apenas
    de prever com maior precisão o momento ideal para realizar a manutenção, mas também
    de prescrever intervenções específicas baseadas em dados, maximizando a vida útil dos
    componentes sem comprometer a confiabilidade do sistema produtivo.
    """
    story.append(Paragraph(intro_text, body_style))
    story.append(PageBreak())

    # === 3. DESCRIÇÃO DO MÉTODO ===
    story.append(Paragraph("3. Descrição do Método Utilizado", heading1_style))
    metodo_text = """
    A metodologia seguiu um fluxo iterativo de cinco etapas principais, garantindo que o
    pipeline funcione de maneira eficiente e seja avaliado continuamente:
    """
    story.append(Paragraph(metodo_text, body_style))

    # Sub-seções do método
    story.append(Paragraph("Coleta e Integração", heading2_style))
    story.append(Paragraph(
        "Nesta etapa, os arquivos CSV provenientes de diferentes injetoras são reunidos. "
        "Cada arquivo traz dados como quantidade produzida, quantidade refugada, datas de "
        "produção, códigos de produto e informações sobre o consumo de massa. O objetivo é "
        "unificar esses dados em um DataFrame único para viabilizar análises consolidadas.",
        body_style
    ))

    story.append(Paragraph("Pré-processamento e Limpeza", heading2_style))
    story.append(Paragraph(
        "Uma vez que os dados estejam integrados, procede-se à limpeza. Isso envolve a "
        "remoção de possíveis duplicatas, conversão de campos de data para o formato "
        "apropriado (datetime), tratamento de valores ausentes, padronização de colunas e "
        "exclusão de registros inconsistentes. Além disso, nesta etapa ocorrem transformações "
        "como a codificação One-Hot para variáveis categóricas.",
        body_style
    ))

    story.append(Paragraph("Análise Exploratória", heading2_style))
    story.append(Paragraph(
        "Com os dados preparados, explora-se o conjunto para identificar padrões, outliers e "
        "tendências. São calculadas estatísticas descritivas (média, mediana, desvio padrão) "
        "e produzidos gráficos como histogramas, boxplots e gráficos de dispersão.",
        body_style
    ))

    story.append(Paragraph("Modelagem (Treinamento)", heading2_style))
    story.append(Paragraph(
        "Na fase de modelagem, definimos a variável-alvo (tempo até a próxima manutenção) "
        "e selecionamos algoritmos de aprendizado de máquina. Modelos como Regressão Linear, "
        "Decision Tree, Random Forest e XGBoost são treinados nos dados de treinamento "
        "(80% do total).",
        body_style
    ))

    story.append(Paragraph("Validação e Avaliação", heading2_style))
    story.append(Paragraph(
        "Após o treinamento, avalia-se o desempenho dos modelos em dados de teste (20%). "
        "Métricas como R², MSE e MAE são calculadas para verificar o quão próximo cada "
        "modelo chega do valor real.",
        body_style
    ))
    story.append(PageBreak())

    # === 4. FERRAMENTAS UTILIZADAS ===
    story.append(Paragraph("4. Ferramentas Utilizadas", heading1_style))
    ferramentas = [
        ("Python (3.x)", "para toda a etapa de análise e modelagem"),
        ("Pandas e NumPy", "para manipulação de dados e transformações"),
        ("Scikit-learn", "como principal biblioteca de Machine Learning"),
        ("XGBoost", "para algoritmos de gradient boosting"),
        ("Matplotlib/Seaborn", "para visualizações de gráficos"),
        ("Jupyter Notebooks", "para execução interativa de código"),
    ]
    for ferramenta, desc in ferramentas:
        story.append(Paragraph(f"<b>{ferramenta}</b> {desc}.", body_style))
    story.append(PageBreak())

    # === 5. MODELOS UTILIZADOS ===
    story.append(Paragraph("5. Modelos Estocásticos e Estatísticos Utilizados", heading1_style))

    story.append(Paragraph("5.1 Regressão Linear", heading2_style))
    story.append(Paragraph(
        "A Regressão Linear é uma técnica estatística clássica que busca entender a relação "
        "entre uma variável dependente (tempo até a manutenção) e variáveis independentes. "
        "O modelo assume a forma: y ≈ β0 + β1X1 + β2X2 + ... + βnXn.",
        body_style
    ))

    story.append(Paragraph("5.2 Árvores de Decisão (Decision Tree)", heading2_style))
    story.append(Paragraph(
        "As Árvores de Decisão organizam os dados em subdivisões baseadas em critérios "
        "estatísticos, como o erro quadrático médio. São úteis para dados com interações "
        "complexas ou não lineares.",
        body_style
    ))

    story.append(Paragraph("5.3 Random Forest", heading2_style))
    story.append(Paragraph(
        "O Random Forest combina diversas Árvores de Decisão para tornar as previsões mais "
        "precisas e estáveis. Geralmente apresenta menor variância que uma única árvore, "
        "sendo mais robusto a ruídos e outliers.",
        body_style
    ))

    story.append(Paragraph("5.4 XGBoost (Extreme Gradient Boosting)", heading2_style))
    story.append(Paragraph(
        "O XGBoost é um modelo avançado que constrói árvores de decisão sequencialmente, "
        "onde cada nova árvore aprende com os erros da anterior. Utiliza técnicas de "
        "regularização (L1 e L2) para evitar overfitting.",
        body_style
    ))
    story.append(PageBreak())

    # === 6. MÉTRICAS ===
    story.append(Paragraph("6. Métricas", heading1_style))

    story.append(Paragraph("6.1 Erro Quadrático Médio (EQM/MSE)", heading2_style))
    story.append(Paragraph(
        "Mede a diferença entre os valores previstos pelo modelo e os valores reais, "
        "elevando essas diferenças ao quadrado. Quanto menor o EQM, mais precisas são "
        "as previsões do modelo.",
        body_style
    ))

    story.append(Paragraph("6.2 R-quadrado (Coeficiente de Determinação)", heading2_style))
    story.append(Paragraph(
        "O R-quadrado indica a proporção da variabilidade da variável dependente que é "
        "explicada pelo modelo. Seu valor varia entre 0 e 1, onde valores próximos de 1 "
        "indicam alta qualidade preditiva.",
        body_style
    ))

    story.append(Paragraph("6.3 Erro Absoluto Médio (EAM/MAE)", heading2_style))
    story.append(Paragraph(
        "Mede a média das diferenças absolutas entre os valores previstos e os reais. "
        "Ao contrário do EQM, não é tão sensível a erros extremos.",
        body_style
    ))
    story.append(PageBreak())

    # === 7. DICIONÁRIO COMPLETO DE VARIÁVEIS ===
    story.append(Paragraph("7. Dicionário Completo de Variáveis do Pipeline", heading1_style))

    story.append(Paragraph(
        "Esta seção documenta todas as variáveis utilizadas no pipeline de Machine Learning, "
        "desde a coleta de dados brutos até as features utilizadas na modelagem. Cada variável "
        "é descrita com seu tipo, origem, uso e exemplos.",
        body_style
    ))

    features = results.get("data", {}).get("features", [])
    total_features = len(features) if features else 0

    story.append(Paragraph(
        f"<b>Total de variáveis no modelo final:</b> {total_features} features",
        body_style
    ))
    story.append(PageBreak())

    # --- 7.1 VARIÁVEIS ORIGINAIS (DADOS BRUTOS) ---
    story.append(Paragraph("7.1 Variáveis Originais (Dados Brutos)", heading1_style))
    story.append(Paragraph(
        f"<b>{VARIAVEIS_PIPELINE['dados_brutos']['descricao']}</b>",
        body_style
    ))

    # Tabela de variáveis originais
    for var in VARIAVEIS_PIPELINE['dados_brutos']['variaveis']:
        story.append(Paragraph(f"<b>{var['nome']}</b>", heading2_style))

        var_info = f"""
        <b>Nome Processado:</b> {var['nome_processado']}<br/>
        <b>Tipo:</b> {var['tipo']}<br/>
        <b>Descrição:</b> {var['descricao']}<br/>
        <b>Exemplo:</b> {var['exemplo']}<br/>
        <b>Origem:</b> {var['origem']}<br/>
        <b>Uso no Pipeline:</b> {var['uso']}
        """
        if 'unidade' in var:
            var_info = var_info.replace('<b>Exemplo:</b>', f"<b>Unidade:</b> {var['unidade']}<br/><b>Exemplo:</b>")

        story.append(Paragraph(var_info, body_style))
        story.append(Spacer(1, 0.3*cm))

    story.append(PageBreak())

    # --- 7.2 VARIÁVEIS GERADAS NO PRÉ-PROCESSAMENTO ---
    story.append(Paragraph("7.2 Variáveis Geradas no Pré-processamento (Etapa 2)", heading1_style))
    story.append(Paragraph(
        f"<b>{VARIAVEIS_PIPELINE['preprocessamento']['descricao']}</b>",
        body_style
    ))

    for var in VARIAVEIS_PIPELINE['preprocessamento']['variaveis']:
        story.append(Paragraph(f"<b>{var['nome']}</b>", heading2_style))

        var_info = f"""
        <b>Nome Processado:</b> {var['nome_processado']}<br/>
        <b>Tipo:</b> {var['tipo']}<br/>
        <b>Descrição:</b> {var['descricao']}<br/>
        <b>Unidade:</b> {var.get('unidade', 'N/A')}<br/>
        <b>Exemplo:</b> {var['exemplo']}<br/>
        <b>Origem:</b> {var['origem']}<br/>
        <b>Uso no Pipeline:</b> {var['uso']}
        """
        if 'formula' in var:
            var_info += f"<br/><b>Fórmula:</b> <i>{var['formula']}</i>"

        story.append(Paragraph(var_info, body_style))
        story.append(Spacer(1, 0.3*cm))

    story.append(PageBreak())

    # --- 7.3 VARIÁVEIS DE ENCODING (ONE-HOT) ---
    story.append(Paragraph("7.3 Variáveis de Encoding (One-Hot)", heading1_style))
    story.append(Paragraph(
        f"<b>{VARIAVEIS_PIPELINE['encoding']['descricao']}</b>",
        body_style
    ))

    for var in VARIAVEIS_PIPELINE['encoding']['variaveis']:
        story.append(Paragraph(f"<b>{var['nome']}</b>", heading2_style))

        var_info = f"""
        <b>Nome Processado:</b> {var['nome_processado']}<br/>
        <b>Tipo:</b> {var['tipo']}<br/>
        <b>Descrição:</b> {var['descricao']}<br/>
        <b>Exemplo:</b> {var['exemplo']}<br/>
        <b>Origem:</b> {var['origem']}<br/>
        <b>Uso no Pipeline:</b> {var['uso']}<br/>
        <b>Categorias:</b> <i>{var.get('categorias', 'Variável')}</i>
        """

        story.append(Paragraph(var_info, body_style))
        story.append(Spacer(1, 0.3*cm))

    story.append(PageBreak())

    # --- 7.4 HISTÓRICO DE TROCA DE PEÇAS POR EQUIPAMENTO ---
    story.append(Paragraph("7.4 Histórico de Troca de Peças por Equipamento", heading1_style))
    story.append(Paragraph(
        f"<b>{VARIAVEIS_PIPELINE['troca_pecas_equipamentos']['descricao']}</b>",
        body_style
    ))

    # Criar tabela de equipamentos com dados de troca de peças
    # CARREGAMENTO DINÂMICO de equipment_stats.json
    equip_data = load_equipment_stats()
    if not equip_data:
        equip_data = VARIAVEIS_PIPELINE.get('troca_pecas_equipamentos', {}).get('equipamentos', {})

    # Tabela: #, Equipamento, Penúltima Troca, Última Troca, Intervalo, Média/Mín/Máx por equipamento
    table_data = [[
        "#", "Equipamento", "Penúltima Troca", "Última Troca",
        "Intervalo (dias)", "Média Hist. (dias)", "Mín (dias)", "Máx (dias)"
    ]]

    for idx, (equip, dados) in enumerate(sorted(equip_data.items()), start=1):
        if isinstance(dados, dict):
            ultima = dados.get('ultima_troca') or dados.get('data_ultima_manutencao') or 'N/A'
            penultima = dados.get('data_penultima_manutencao') or 'N/A'
            intervalo = dados.get('intervalo_dias') or dados.get('intervalo_manutencao_dias') or 'N/A'
            media = dados.get('media_dias_manutencao')
            min_d = dados.get('min_dias_manutencao')
            max_d = dados.get('max_dias_manutencao')
            media_str = f"{media:.0f}" if isinstance(media, (int, float)) else 'N/A'
            min_str = f"{min_d:.0f}" if isinstance(min_d, (int, float)) else 'N/A'
            max_str = f"{max_d:.0f}" if isinstance(max_d, (int, float)) else 'N/A'
        else:
            ultima = dados
            penultima = 'N/A'
            intervalo = 'N/A'
            media_str = min_str = max_str = 'N/A'
        table_data.append([
            str(idx), equip, str(penultima), str(ultima),
            str(intervalo), media_str, min_str, max_str
        ])

    equip_table = Table(
        table_data,
        colWidths=[1*cm, 2.5*cm, 2.5*cm, 2.5*cm, 2.2*cm, 2.5*cm, 1.8*cm, 1.8*cm]
    )
    equip_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#F77F00')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 8),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
        ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#F5F5F5')),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.black),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
    ]))
    story.append(equip_table)
    story.append(Spacer(1, 0.5*cm))

    # --- 7.5 RESUMO DAS FEATURES UTILIZADAS NO MODELO ---
    story.append(Paragraph("7.5 Features Utilizadas no Modelo Final", heading1_style))

    if features:
        story.append(Paragraph(
            f"O modelo final utiliza <b>{len(features)}</b> features numéricas. "
            "Abaixo está a lista completa das features após todo o pré-processamento:",
            body_style
        ))

        # Agrupar features por categoria
        features_originais = [f for f in features if not any(f.startswith(p) for p in ['Equipamento_', 'Cod_Produto_', 'Descricao_da_massa_', 'Cod_Un_']) and 'Acumulado' not in f]
        features_acumuladas = [f for f in features if 'Acumulado' in f]
        features_encoding = [f for f in features if any(f.startswith(p) for p in ['Equipamento_', 'Cod_Produto_', 'Descricao_da_massa_', 'Cod_Un_'])]

        if features_originais:
            story.append(Paragraph(f"<b>Features Numéricas Originais ({len(features_originais)}):</b>", body_style))
            for f in features_originais:
                story.append(Paragraph(f"• {f}", styles['Normal']))
            story.append(Spacer(1, 0.3*cm))

        if features_acumuladas:
            story.append(Paragraph(f"<b>Features Acumuladas ({len(features_acumuladas)}):</b>", body_style))
            for f in features_acumuladas:
                story.append(Paragraph(f"• {f}", styles['Normal']))
            story.append(Spacer(1, 0.3*cm))

        if features_encoding:
            story.append(Paragraph(f"<b>Features de Encoding ({len(features_encoding)}):</b>", body_style))
            # Mostrar apenas algumas para não poluir
            if len(features_encoding) > 15:
                story.append(Paragraph(", ".join(features_encoding[:15]) + f"... (e mais {len(features_encoding)-15} features)", styles['Normal']))
            else:
                story.append(Paragraph(", ".join(features_encoding), styles['Normal']))

    story.append(PageBreak())

    # === 8. DADOS REMOVIDOS ===
    story.append(Paragraph("8. Dados Removidos da Análise", heading1_style))
    story.append(Paragraph(
        "Durante a fase de pré-processamento, foram eliminados ou desconsiderados alguns "
        "registros e variáveis por apresentarem inconsistências:",
        body_style
    ))
    dados_removidos = [
        "Registros Duplicados: excluídos para evitar distorções na contagem.",
        "Linhas com Dados Inconsistentes: valores impossíveis ou datas inválidas.",
        "Valores Nulos Irrecuperáveis: quando variáveis essenciais estavam ausentes.",
        "Variáveis Redundantes: que não acrescentavam novas informações.",
    ]
    for item in dados_removidos:
        story.append(Paragraph(f"• {item}", body_style))

    # Nota de integridade: diferenças entre versões do relatório
    story.append(Paragraph(
        "<b>Nota sobre Integridade dos Dados (aplicável a partir desta versão):</b>",
        body_style
    ))
    story.append(Paragraph(
        "Comparando esta versão do relatório com versões anteriores (ex.: R14, R15), "
        "podem haver divergências significativas no total de registros e na contagem de "
        "variáveis. Os números anteriores estavam afetados por três problemas de ingestão "
        "que foram identificados e corrigidos:",
        body_style
    ))
    correcoes = [
        "<b>Dupla contagem de planilhas de produção:</b> quando o arquivo DadosProducao*.xlsx "
        "ficava acidentalmente em data/raw/, o pipeline carregava os mesmos registros duas "
        "vezes (uma via IJ-*.csv já splitado pelo s00, outra via leitura direta do xlsx), "
        "inflando o total. Corrigido movendo o xlsx para data/arquivo_unico_processado/ após "
        "o split, e o carregador direto só vê arquivos ainda não processados.",
        "<b>Sobrescrita silenciosa do baseline histórico:</b> um bloco except no s00 "
        "substituía os CSVs existentes por apenas os registros novos em caso de qualquer "
        "erro de leitura. A ingestão passou a preservar o histórico e logar o erro em vez "
        "de descartar dados.",
        "<b>Interpretação errada de datas em formato ISO:</b> aplicar dayfirst=True sobre "
        "datas no padrão YYYY-MM-DD (como \"2026-01-12\") fazia o pandas reinterpretá-las "
        "trocando mês e dia (resultando em \"2026-12-01\" — 1° de dezembro no futuro). "
        "A leitura agora tenta ISO 8601 primeiro e só recorre a dayfirst=True como fallback "
        "para valores no formato brasileiro dd/mm/yyyy.",
    ]
    for item in correcoes:
        story.append(Paragraph(f"• {item}", body_style))
    story.append(Paragraph(
        "Como consequência, este relatório apresenta um conjunto de dados menor em volume "
        "mas íntegro em conteúdo, com datas cronologicamente corretas e sem duplicações "
        "artificiais. As métricas de modelo permanecem comparáveis em qualidade (R² na "
        "mesma faixa), confirmando que o dataset depurado preserva o sinal preditivo.",
        body_style
    ))
    story.append(PageBreak())

    # === 9. DESEMPENHO DE MODELOS ===
    story.append(Paragraph("9. Desempenho de Modelos", heading1_style))

    story.append(Paragraph(
        f"O modelo {model_name.upper()} foi selecionado como o melhor modelo nos experimentos, "
        "destacando a capacidade de capturar padrões complexos de desgaste.",
        body_style
    ))

    # Tabela de resultados
    if metrics:
        story.append(Paragraph("Quadro 1: Desempenho do Modelo Selecionado:", body_style))

        table_data = [
            ["Modelo", "EQM (MSE)", "R-Quadrado", "EAM (MAE)"],
            [model_name.upper(), f"{mse:.2f}", f"{r2:.4f}", f"{mae:.2f}"]
        ]

        table = Table(table_data, colWidths=[4*cm, 3*cm, 3*cm, 3*cm])
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#F77F00')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 11),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#F5F5F5')),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ]))
        story.append(table)
    story.append(PageBreak())

    # === 8b. GRÁFICOS DE ANÁLISE ===
    story.append(Paragraph("Gráficos de Análise", heading1_style))
    story.append(Paragraph(
        "Os gráficos abaixo apresentam as análises exploratórias realizadas durante o "
        "desenvolvimento do modelo, incluindo matrizes de correlação, análises de distribuição "
        "e visualizações das principais features.",
        body_style
    ))

    # Diretórios onde procurar gráficos (em ordem de prioridade)
    # 1. Diretório local (onde o pipeline é executado)
    # 2. Diretório scripts/eda_plots
    # 3. Diretórios de análises exploratórias existentes (fallback)
    script_dir = Path(__file__).parent.parent  # sabo/sabo/scripts
    plots_directories = [
        Path("eda_plots"),  # Diretório relativo ao CWD (scripts/)
        script_dir / "eda_plots",  # Caminho absoluto para scripts/eda_plots
        Path(__file__).parent.parent / "eda_plots",  # Alternativo
        script_dir.parent / "exploratory_analise",  # Gráficos existentes (fallback)
        script_dir.parent / "exploratory_dado_proposto" / "classification_images",
    ]

    # Log dos diretórios sendo verificados
    print("  Procurando gráficos em:")
    for d in plots_directories:
        exists = "✓" if d.exists() else "✗"
        print(f"    {exists} {d}")

    # Mapeamento de gráficos prioritários com legendas detalhadas
    # Formato: (nome_arquivo, legenda, base_calculo)
    priority_plots = [
        ("correlation_matrix.png", f"Matriz de Correlação entre Variáveis - Período: {periodo_inicio} a {periodo_fim}", "Base: Quantidade de peças"),
        ("correlation_matrix_full.png", f"Matriz de Correlação Completa - Período: {periodo_inicio} a {periodo_fim}", "Base: Quantidade de peças"),
        ("heatmap_correlacao.png", f"Heatmap de Correlação - Período: {periodo_inicio} a {periodo_fim}", "Base: Quantidade de peças"),
        ("correlation_matrix_heatmap.png", f"Heatmap de Correlação - Período: {periodo_inicio} a {periodo_fim}", "Base: Quantidade de peças"),
        ("consumo_massa_vs_qtd_produzida.png", f"Consumo de Massa vs Quantidade Produzida - Período: {periodo_inicio} a {periodo_fim}", "Base: Quantidade de peças (unidades)"),
        ("consumo_vs_producao.png", f"Análise de Consumo vs Produção - Período: {periodo_inicio} a {periodo_fim}", "Base: Quantidade de peças"),
        ("boxplot_consumo_massa_total.png", f"Boxplot do Consumo de Massa Total - Período: {periodo_inicio} a {periodo_fim}", "Base: Quantidade (kg)"),
        ("ano_construcao_vs_consumo_massa.png", f"Ano de Construção vs Consumo de Massa - Período: {periodo_inicio} a {periodo_fim}", "Base: Anos / Quantidade (kg)"),
        ("analise_temporal.png", f"Análise Temporal - Período: {periodo_inicio} a {periodo_fim} ({total_dias} dias)", "Base: Dias"),
        ("matriz_urgencia.png", f"Distribuição de Urgência de Manutenção - Período: {periodo_inicio} a {periodo_fim}", "Base: Classificação por nível de urgência"),
        ("scatter_plots.png", f"Gráficos de Dispersão - Período: {periodo_inicio} a {periodo_fim}", "Base: Quantidade de peças"),
        ("scatter_plots_features.png", f"Scatter Plots das Features por Equipamento - Período: {periodo_inicio} a {periodo_fim}", "Base: Valores das variáveis coloridos por equipamento"),
        ("resumo_equipamentos.png", f"Resumo por Equipamento - Período: {periodo_inicio} a {periodo_fim}", "Base: Produção, Refugo e Dias até Manutenção"),
        ("dispersao_target.png", f"Dispersão do Target - Período: {periodo_inicio} a {periodo_fim}", "Base: Valores preditos vs reais"),
        ("histogramas.png", f"Histogramas das Variáveis - Período: {periodo_inicio} a {periodo_fim}", "Base: Frequência de ocorrência"),
        ("boxplots.png", f"Boxplots das Variáveis - Período: {periodo_inicio} a {periodo_fim}", "Base: Distribuição estatística"),
    ]

    plots_added = 0
    max_plots = 8  # Aumentado para incluir gráficos por equipamento
    added_names = set()  # Evitar duplicatas

    for plot_info in priority_plots:
        if plots_added >= max_plots:
            break

        plot_name, caption, base_calculo = plot_info

        # Verificar nome base para evitar duplicatas similares
        base_name = plot_name.replace("_full", "").replace("_heatmap", "")
        if base_name in added_names:
            continue

        # Procurar o gráfico em todos os diretórios
        for plots_dir in plots_directories:
            if not plots_dir.exists():
                continue

            plot_path = plots_dir / plot_name
            if plot_path.exists():
                try:
                    # Calcular dimensões proporcionais
                    img_width = 16*cm
                    img_height = 11*cm

                    # Adicionar gráfico ao relatório
                    img = Image(str(plot_path), width=img_width, height=img_height)
                    story.append(img)
                    story.append(Spacer(1, 0.3*cm))

                    # Legenda do gráfico com período e base de cálculo
                    story.append(Paragraph(
                        f"<i>Figura {plots_added + 1}: {caption}</i>",
                        ParagraphStyle('Caption', parent=styles['Normal'], fontSize=10, alignment=TA_CENTER)
                    ))
                    story.append(Paragraph(
                        f"<i>{base_calculo}</i>",
                        ParagraphStyle('CaptionBase', parent=styles['Normal'], fontSize=9, alignment=TA_CENTER, textColor=colors.gray)
                    ))
                    story.append(Spacer(1, 0.8*cm))

                    plots_added += 1
                    added_names.add(base_name)
                    print(f"  ✓ Gráfico adicionado: {plot_name}")
                    break  # Encontrou o gráfico, não procurar mais
                except Exception as e:
                    print(f"  ⚠ Erro ao adicionar gráfico {plot_name}: {e}")

    if plots_added == 0:
        story.append(Paragraph(
            "<i>Gráficos não disponíveis. Execute a etapa 3 (EDA) ou 3b (Análises Avançadas) para gerá-los.</i>",
            styles['Normal']
        ))
    else:
        print(f"  Total de gráficos adicionados ao relatório: {plots_added}")

    # === 10. PERGUNTAS E RESPOSTAS ===
    # Sem PageBreak para evitar página vazia após gráficos
    story.append(Spacer(1, 1*cm))
    story.append(Paragraph("10. Perguntas e Respostas", heading1_style))

    story.append(Paragraph("10.1 Como interpretar os gráficos?", heading2_style))
    story.append(Paragraph(
        f"Os gráficos apresentados utilizam dados do período de <b>{periodo_inicio}</b> a <b>{periodo_fim}</b>. "
        "A base de cálculo varia conforme o tipo de análise: gráficos temporais utilizam dias ou meses "
        "como unidade, enquanto gráficos de produção utilizam quantidade de peças. A legenda de cada "
        "gráfico indica a unidade utilizada.",
        body_style
    ))

    story.append(Paragraph("10.2 Qual a confiabilidade das previsões?", heading2_style))
    story.append(Paragraph(
        f"O modelo {model_name.upper()} apresentou R² de {r2:.4f}, o que significa que aproximadamente "
        f"{r2*100:.1f}% da variação nos dados é explicada pelo modelo. O erro médio absoluto (MAE) de "
        f"{mae:.2f} indica a margem de erro típica das previsões.",
        body_style
    ))

    story.append(Paragraph("10.3 Os dados são representativos?", heading2_style))
    story.append(Paragraph(
        f"A análise foi baseada em {data_shape[0]} registros coletados ao longo de {total_dias} dias "
        f"({total_dias // 30} meses). Para maior robustez estatística, recomenda-se mínimo de 30 "
        "observações por cenário, sendo ideal entre 50 e 100.",
        body_style
    ))

    story.append(Paragraph("10.4 Como o modelo foi validado?", heading2_style))
    story.append(Paragraph(
        f"Os dados foram divididos em conjunto de treinamento ({train_size} amostras, 80%) e teste "
        f"({test_size} amostras, 20%). O modelo foi treinado apenas com dados de treinamento e "
        "avaliado em dados nunca vistos (teste), garantindo avaliação imparcial.",
        body_style
    ))
    story.append(PageBreak())

    # === 11. PREVISÃO DE TROCA DE PEÇAS ===
    story.append(Paragraph("11. Previsão de Troca de Peças", heading1_style))

    # Carregar dados de equipamentos
    today = datetime.now()
    equip_data = load_equipment_stats()
    if not equip_data:
        equip_data = VARIAVEIS_PIPELINE.get('troca_pecas_equipamentos', {}).get('equipamentos', {})

    story.append(Paragraph(
        "Esta seção consolida, em um único arquivo CSV, a previsão de manutenção por equipamento "
        "combinando dados históricos, agendamento SAP (preventivas RM.195) e prescrição do modelo "
        "de Machine Learning. As colunas do arquivo são: Equipamento, Última Data de Produção, "
        "Última Manutenção, Intervalo entre Penúltima e Última Manutenção (dias), Data Próxima "
        "Manutenção Agendada SAP, Dias da Última Manutenção até a Data SAP, Diferença entre Data "
        "SAP e Hoje, Data Próxima Manutenção Prescrita pelo Machine Learning e Diferença em Dias "
        "entre Data ML e Data SAP (sinal positivo: ML após SAP; negativo: ML antes do SAP).",
        body_style
    ))
    story.append(Spacer(1, 0.3*cm))

    # Gerar CSV consolidado
    csv_path = generate_previsao_manutencao_csv(equip_data)
    csv_abs = Path(csv_path).resolve()

    # Gerar PPTX consolidado
    pptx_abs = None
    try:
        pptx_path = merge_componentes_pptx()
        if pptx_path:
            pptx_abs = Path(pptx_path).resolve()
    except Exception as e:
        print(f"  ⚠ Falha na consolidação de PPTX: {e}")

    story.append(Paragraph(
        f"• <b>Arquivo CSV:</b> "
        f"<link href='file://{csv_abs}' color='blue'><u>{csv_abs.name}</u></link>",
        body_style
    ))
    if pptx_abs is not None:
        story.append(Paragraph(
            f"• <b>Apresentação consolidada (PPTX):</b> "
            f"<link href='file://{pptx_abs}' color='blue'><u>{pptx_abs.name}</u></link>",
            body_style
        ))

    story.append(PageBreak())

    # NOTA: subitens 11.1 a 11.4 removidos — substituídos pelo CSV consolidado acima.

    # === 12. ANÁLISE MENSAL POR EQUIPAMENTO ===
    story.append(Paragraph("12. Análise Mensal por Equipamento", heading1_style))

    monthly_charts = results.get("monthly_charts", [])
    if monthly_charts:
        story.append(Paragraph(
            f"Os gráficos abaixo apresentam a análise mensal de cada equipamento, incluindo "
            f"quantidade de peças produzidas, peças refugadas e momentos de manutenção. "
            f"Período exibido: últimos 6 meses dos dados analisados ({periodo_inicio} a {periodo_fim}).",
            body_style
        ))
        story.append(Spacer(1, 0.5*cm))

        import math
        for idx, chart_info in enumerate(monthly_charts):
            # chart_info pode ser (path, n_charts) ou apenas path
            if isinstance(chart_info, (list, tuple)):
                chart_path, n_charts_on_page = chart_info
            else:
                chart_path, n_charts_on_page = chart_info, 8

            if Path(chart_path).exists():
                try:
                    # Altura proporcional ao número de linhas usadas
                    n_rows = math.ceil(n_charts_on_page / 2)
                    img_height = 23 * n_rows / 4  # 4 linhas = 23cm
                    img = Image(str(chart_path), width=17*cm, height=img_height*cm)
                    story.append(img)
                    story.append(Spacer(1, 0.3*cm))
                    story.append(Paragraph(
                        f"<i>Figura: Análise Mensal - Página {idx + 1}</i>",
                        ParagraphStyle('Caption', parent=styles['Normal'], fontSize=10, alignment=TA_CENTER)
                    ))
                    if idx < len(monthly_charts) - 1:
                        story.append(PageBreak())
                except Exception as e:
                    print(f"  ⚠ Erro ao incluir gráfico mensal {chart_path}: {e}")
    else:
        story.append(Paragraph(
            "Gráficos mensais não disponíveis. Verifique se data_raw.csv existe e contém dados suficientes.",
            body_style
        ))

    story.append(PageBreak())

    # === 13. CONSIDERAÇÕES ===
    story.append(Paragraph("13. Considerações", heading1_style))
    story.append(Paragraph(
        "O desenvolvimento de um projeto de manutenção prescritiva para injetoras de borracha "
        "demonstrou o potencial de algoritmos de Machine Learning não apenas para prever o "
        "momento ideal de manutenção, mas para prescrever intervenções específicas baseadas "
        "nos padrões identificados.",
        body_style
    ))
    story.append(Paragraph(
        "O ensaio de diferentes algoritmos (Regressão Linear, Decision Tree, Random Forest e "
        "XGBoost) evidenciou que modelos de ensemble tendem a produzir resultados superiores "
        "quando há volume de dados adequado.",
        body_style
    ))
    story.append(Spacer(1, 0.8*cm))

    # === 14. PRÓXIMOS PASSOS ===
    story.append(Paragraph("14. Próximos Passos", heading1_style))
    proximos_passos = [
        ("Shadow Teste", "Executar o modelo em operação paralela à linha produtiva, sem intervenções reais."),
        ("Sistema de Recomendações", "Criar regras que traduzam resultados em recomendações de intervenção."),
        ("Implantação em Produção", "Adoção gradativa começando por linha-piloto."),
        ("Validação Contínua", "Incorporar dados recentes de forma recorrente."),
        ("Incorporar Variáveis de Campo", "Incluir dados de temperatura, vibração e pressão."),
    ]
    for titulo, desc in proximos_passos:
        story.append(Paragraph(f"<b>{titulo}:</b> {desc}", body_style))
    story.append(Spacer(1, 0.8*cm))

    # === 15. RECOMENDAÇÕES ===
    story.append(Paragraph("15. Recomendações", heading1_style))
    recomendacoes = [
        "Desenvolvimento de Biblioteca de Recomendações Prescritivas",
        "Plano Estruturado de Coleta e Qualidade de Dados",
        "Cultura de manutenção baseada em dados",
        "Integração com Sistemas de Produção",
        "Pipelines de Teste e Retreinamento",
        "Explorar Novos Modelos e Técnicas de Interpretação",
        "Foco nos Benefícios Financeiros e Estratégicos",
    ]
    for rec in recomendacoes:
        story.append(Paragraph(f"• {rec}", body_style))
    story.append(PageBreak())

    # === 16. CONSIDERAÇÕES FINAIS ===
    story.append(Paragraph("16. Considerações Finais", heading1_style))
    story.append(Paragraph(
        f"Os resultados preliminares revelaram boas perspectivas para métodos como "
        f"{model_name.upper()}. A incorporação de novas amostras tende a aumentar a precisão "
        "das previsões, viabilizando uma manutenção realmente prescritiva, que não apenas "
        "prevê quando intervir, mas também determina as ações específicas necessárias.",
        body_style
    ))
    story.append(Paragraph(
        "Como parâmetro de amostragem mínima, recomenda-se reunir ao menos 30 observações "
        "para cada cenário relevante, sendo a faixa ideal entre 50 e 100 observações para "
        "garantir maior robustez estatística.",
        body_style
    ))
    story.append(PageBreak())

    # === 17. CONSUMO DE MASSA POR EQUIPAMENTO E COMPOSTO ===
    story.append(Paragraph("17. Consumo de Massa por Equipamento e Composto", heading1_style))
    story.append(Paragraph(
        "Esta seção detalha, para cada equipamento, a quantidade de cada tipo de composto "
        "(massa de borracha) utilizado no período analisado, bem como o volume em "
        "quilogramas consumidos. A massa consumida é calculada como "
        "<b>Kg = Qtd. Produzida × Consumo (Kg/100pçs) / 100</b>. O objetivo é apoiar "
        "decisões de planejamento de insumos, identificar equipamentos com maior consumo "
        "de compostos específicos e correlacionar tipo de composto com padrões de desgaste.",
        body_style
    ))
    story.append(Spacer(1, 0.4*cm))

    # 17.1 Heatmap geral
    compound_heatmap = results.get("compound_heatmap")
    if compound_heatmap and Path(compound_heatmap).exists():
        story.append(Paragraph("17.1 Visão Geral — Heatmap Equipamento × Composto", heading2_style))
        story.append(Paragraph(
            "O heatmap abaixo apresenta de forma consolidada o consumo (em kg) de cada "
            "composto por equipamento. Células mais escuras indicam maior consumo. "
            "Compostos estão ordenados da esquerda para a direita em ordem decrescente "
            "de consumo total.",
            body_style
        ))
        try:
            img = Image(str(compound_heatmap), width=17*cm, height=11*cm)
            story.append(img)
            story.append(Paragraph(
                "<i>Figura: Heatmap de consumo de massa (kg) por equipamento e composto.</i>",
                ParagraphStyle('Caption', parent=styles['Normal'], fontSize=10, alignment=TA_CENTER)
            ))
        except Exception as e:
            print(f"  ⚠ Erro ao incluir heatmap de composto: {e}")
        story.append(PageBreak())

    # 17.2 Gráficos individuais por equipamento
    compound_charts = results.get("compound_charts", [])
    if compound_charts:
        story.append(Paragraph("17.2 Consumo Detalhado por Equipamento", heading2_style))
        story.append(Paragraph(
            "Os gráficos a seguir mostram, para cada equipamento, o consumo total de cada "
            "composto em kg, com a quantidade de peças produzidas indicada ao lado de cada "
            "barra. Os equipamentos são apresentados em ordem alfabética, 8 por página.",
            body_style
        ))
        story.append(Spacer(1, 0.3*cm))

        import math
        for idx, chart_info in enumerate(compound_charts):
            if isinstance(chart_info, (list, tuple)):
                chart_path, n_charts_on_page = chart_info
            else:
                chart_path, n_charts_on_page = chart_info, 8

            if Path(chart_path).exists():
                try:
                    n_rows = math.ceil(n_charts_on_page / 2)
                    img_height = 23 * n_rows / 4
                    img = Image(str(chart_path), width=17*cm, height=img_height*cm)
                    story.append(img)
                    story.append(Spacer(1, 0.3*cm))
                    story.append(Paragraph(
                        f"<i>Figura: Consumo por Composto — Página {idx + 1}</i>",
                        ParagraphStyle('Caption', parent=styles['Normal'], fontSize=10, alignment=TA_CENTER)
                    ))
                    if idx < len(compound_charts) - 1:
                        story.append(PageBreak())
                except Exception as e:
                    print(f"  ⚠ Erro ao incluir gráfico de composto {chart_path}: {e}")
        story.append(PageBreak())

    # 17.3 Tabela resumo
    compound_rows = results.get("compound_summary_rows")
    if compound_rows and len(compound_rows) > 1:
        story.append(Paragraph("17.3 Tabela Detalhada — Peças e Kg por Equipamento/Composto", heading2_style))
        story.append(Paragraph(
            "Tabela consolidada com todos os pares equipamento × composto observados no "
            "período. A coluna <b>Peças produzidas</b> corresponde à soma de Qtd. Produzida; "
            "<b>Massa consumida (kg)</b> corresponde à soma de Qtd. Produzida × Consumo/100.",
            body_style
        ))
        story.append(Spacer(1, 0.2*cm))

        # Paginar tabela em blocos para evitar tabelas gigantes únicas
        header = compound_rows[0]
        body_rows = compound_rows[1:]
        chunk_size = 40
        for chunk_idx in range(0, len(body_rows), chunk_size):
            chunk = [header] + body_rows[chunk_idx:chunk_idx + chunk_size]
            table = Table(chunk, colWidths=[2.8*cm, 6.5*cm, 3.5*cm, 4.0*cm], repeatRows=1)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#F77F00')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 9),
                ('FONTSIZE', (0, 1), (-1, -1), 8),
                ('ALIGN', (2, 1), (-1, -1), 'RIGHT'),
                ('ALIGN', (0, 0), (-1, 0), 'CENTER'),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('GRID', (0, 0), (-1, -1), 0.3, colors.grey),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#F8F8F8')]),
            ]))
            story.append(table)
            story.append(Spacer(1, 0.3*cm))
    else:
        story.append(Paragraph(
            "Dados de consumo por composto não disponíveis — verifique se data_raw.csv "
            "contém as colunas 'Descrição da massa (Composto)' e "
            "'Consumo de massa no item em (Kg/100pçs)'.",
            body_style
        ))

    story.append(PageBreak())

    # === REFERÊNCIAS ===
    story.append(Paragraph("Referência Bibliográfica", heading1_style))
    referencias = [
        "Breiman, L. (2001). Random Forests. Machine Learning, 45(1), 5–32.",
        "Breiman, L., Friedman, J. H., Olshen, R. A., & Stone, C. J. (1984). Classification and Regression Trees. Wadsworth.",
        "Chen, T., & Guestrin, C. (2016). XGBoost: A Scalable Tree Boosting System. ACM SIGKDD.",
        "Géron, A. (2019). Hands-On Machine Learning with Scikit-Learn, Keras, and TensorFlow. O'Reilly.",
        "Hastie, T., Tibshirani, R., & Friedman, J. (2009). The Elements of Statistical Learning. Springer.",
    ]
    for ref in referencias:
        story.append(Paragraph(ref, styles['Normal']))
        story.append(Spacer(1, 0.3*cm))

    # Gerar PDF
    doc.build(story)

    return output_path



def generate_text_report(results: dict, output_path: str) -> str:
    """
    Gera relatório em formato texto (fallback quando ReportLab não está disponível).

    Args:
        results: Dados do pipeline
        output_path: Caminho do arquivo

    Returns:
        Caminho do arquivo gerado
    """
    model_name = results.get("model", {}).get("name", "Não identificado")
    metrics = results.get("model", {}).get("metrics", {})
    r2 = metrics.get("r2", 0)
    mse = metrics.get("mse", 0)
    mae = metrics.get("mae", 0)
    data_shape = results.get("data", {}).get("shape", (0, 0))
    features = results.get("data", {}).get("features", [])

    # Gerar seção de variáveis
    variaveis_text = ""

    # Variáveis originais
    variaveis_text += "\n7.1 VARIÁVEIS ORIGINAIS (DADOS BRUTOS)\n"
    variaveis_text += "-" * 50 + "\n"
    for var in VARIAVEIS_PIPELINE['dados_brutos']['variaveis']:
        variaveis_text += f"\n• {var['nome']}\n"
        variaveis_text += f"  Nome Processado: {var['nome_processado']}\n"
        variaveis_text += f"  Tipo: {var['tipo']}\n"
        variaveis_text += f"  Descrição: {var['descricao']}\n"
        variaveis_text += f"  Origem: {var['origem']}\n"
        variaveis_text += f"  Uso: {var['uso']}\n"

    # Variáveis de pré-processamento
    variaveis_text += "\n\n7.2 VARIÁVEIS GERADAS NO PRÉ-PROCESSAMENTO\n"
    variaveis_text += "-" * 50 + "\n"
    for var in VARIAVEIS_PIPELINE['preprocessamento']['variaveis']:
        variaveis_text += f"\n• {var['nome']}\n"
        variaveis_text += f"  Nome Processado: {var['nome_processado']}\n"
        variaveis_text += f"  Tipo: {var['tipo']}\n"
        variaveis_text += f"  Descrição: {var['descricao']}\n"
        variaveis_text += f"  Unidade: {var.get('unidade', 'N/A')}\n"
        variaveis_text += f"  Origem: {var['origem']}\n"
        variaveis_text += f"  Uso: {var['uso']}\n"
        if 'formula' in var:
            variaveis_text += f"  Fórmula: {var['formula']}\n"

    # Variáveis de encoding
    variaveis_text += "\n\n7.3 VARIÁVEIS DE ENCODING (ONE-HOT)\n"
    variaveis_text += "-" * 50 + "\n"
    for var in VARIAVEIS_PIPELINE['encoding']['variaveis']:
        variaveis_text += f"\n• {var['nome']}\n"
        variaveis_text += f"  Tipo: {var['tipo']}\n"
        variaveis_text += f"  Descrição: {var['descricao']}\n"
        variaveis_text += f"  Exemplo: {var['exemplo']}\n"
        variaveis_text += f"  Categorias: {var.get('categorias', 'Variável')}\n"

    # Histórico de troca de peças (carregado dinamicamente)
    variaveis_text += "\n\n7.4 HISTÓRICO DE TROCA DE PEÇAS POR EQUIPAMENTO\n"
    variaveis_text += "-" * 50 + "\n"
    equip_stats = load_equipment_stats()
    if not equip_stats:
        equip_stats = VARIAVEIS_PIPELINE.get('troca_pecas_equipamentos', {}).get('equipamentos', {})
    for equip, data in sorted(equip_stats.items()):
        variaveis_text += f"  {equip}: {data}\n"

    # Features do modelo
    variaveis_text += f"\n\n7.5 FEATURES UTILIZADAS NO MODELO ({len(features)} total)\n"
    variaveis_text += "-" * 50 + "\n"
    if features:
        for i, f in enumerate(features[:30], 1):
            variaveis_text += f"  {i:3d}. {f}\n"
        if len(features) > 30:
            variaveis_text += f"  ... e mais {len(features) - 30} features\n"

    report = f"""
{'=' * 70}
RELATÓRIO DE ENTREGA - SABO
Pipeline de Machine Learning para Manutenção Prescritiva
Data: {datetime.now().strftime('%d/%m/%Y %H:%M')}
{'=' * 70}

SUMÁRIO
-------
1. Resumo
2. Introdução
3. Descrição do Método Utilizado
4. Ferramentas Utilizadas
5. Modelos Utilizados
6. Métricas
7. Dicionário Completo de Variáveis do Pipeline
   7.1 Variáveis Originais (Dados Brutos)
   7.2 Variáveis Geradas no Pré-processamento
   7.3 Variáveis de Encoding (One-Hot)
   7.4 Histórico de Troca de Peças por Equipamento
   7.5 Features Utilizadas no Modelo Final
8. Dados Removidos da Análise
9. Desempenho de Modelos
10. Perguntas e Respostas
11. Previsão de Troca de Peças
12. Análise Mensal por Equipamento
13. Considerações
14. Próximos Passos
15. Recomendações
16. Considerações Finais
17. Consumo de Massa por Equipamento e Composto

{'=' * 70}
RESUMO
{'=' * 70}

Este relatório apresenta os resultados do projeto de prescrição de troca de peças
para extrusoras de borracha modelo Y125.

Dataset: {data_shape[0]} registros, {data_shape[1]} variáveis
Modelo Selecionado: {model_name.upper()}

MÉTRICAS DO MODELO:
  - R² (Coeficiente de Determinação): {r2:.4f}
  - MSE (Erro Quadrático Médio): {mse:.2f}
  - MAE (Erro Absoluto Médio): {mae:.2f}

{'=' * 70}
METODOLOGIA
{'=' * 70}

O pipeline seguiu 5 etapas:
1. Coleta e Integração de Dados
2. Pré-processamento e Limpeza
3. Análise Exploratória (EDA)
4. Modelagem e Treinamento
5. Validação e Avaliação

{'=' * 70}
MODELOS TESTADOS
{'=' * 70}

1. Regressão Linear
2. Decision Tree
3. Random Forest
4. XGBoost

{'=' * 70}
7. DICIONÁRIO COMPLETO DE VARIÁVEIS DO PIPELINE
{'=' * 70}
{variaveis_text}

{'=' * 70}
DESEMPENHO
{'=' * 70}

Modelo Selecionado: {model_name.upper()}
  R²:  {r2:.4f}
  MSE: {mse:.2f}
  MAE: {mae:.2f}

{'=' * 70}
CONSIDERAÇÕES FINAIS
{'=' * 70}

O projeto demonstrou o potencial de algoritmos de Machine Learning para
manutenção prescritiva. Recomenda-se:

- Executar Shadow Teste em ambiente de produção
- Incorporar mais variáveis de campo (temperatura, vibração)
- Ampliar a base de dados para maior robustez estatística
- Implementar retreinamento periódico do modelo

{'=' * 70}
"""

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(report)

    return output_path


def generate_monthly_equipment_charts(inicio=None, fim=None, n_months=6) -> list:
    """
    Gera gráficos mensais por equipamento: produção, refugo e manutenções.

    Args:
        inicio: Data início do período (YYYY-MM-DD), opcional
        fim: Data fim do período (YYYY-MM-DD), opcional
        n_months: Número de meses a exibir (padrão: 6)

    Returns:
        Lista de caminhos dos arquivos PNG gerados
    """
    if not HAS_MATPLOTLIB:
        print("  ⚠ Matplotlib não disponível. Gráficos mensais não serão gerados.")
        return []

    output_dir = Path("eda_plots") / "monthly"
    output_dir.mkdir(parents=True, exist_ok=True)

    # Carregar dados raw
    data_path = Path("data_raw.csv")
    if not data_path.exists():
        print("  ⚠ data_raw.csv não encontrado. Gráficos mensais indisponíveis.")
        return []

    try:
        df = pd.read_csv(data_path)
    except Exception as e:
        print(f"  ⚠ Erro ao carregar data_raw.csv: {e}")
        return []

    # Detectar colunas dinamicamente (suporte a nomes com acentos)
    date_col = next((c for c in df.columns if 'data' in c.lower() and 'produ' in c.lower()), None)
    prod_col = next((c for c in df.columns if 'produzid' in c.lower() and 'acumul' not in c.lower()), None)
    ref_col = next((c for c in df.columns if 'refugad' in c.lower() and 'acumul' not in c.lower()), None)
    retrab_col = next((c for c in df.columns if 'retrabalhad' in c.lower() and 'acumul' not in c.lower()), None)
    equip_col = next((c for c in df.columns if c.lower() == 'equipamento'), None)

    if not all([date_col, prod_col, equip_col]):
        print(f"  ⚠ Colunas necessárias não encontradas. Disponíveis: {list(df.columns)}")
        return []

    # Parse datas
    df[date_col] = pd.to_datetime(df[date_col], errors='coerce')
    df = df.dropna(subset=[date_col])

    if len(df) == 0:
        print("  ⚠ Nenhum registro com data válida.")
        return []

    # Aplicar filtro de período
    if inicio:
        df = df[df[date_col] >= pd.Timestamp(inicio)]
    if fim:
        df = df[df[date_col] <= pd.Timestamp(fim)]

    if len(df) == 0:
        print("  ⚠ Nenhum registro no período especificado.")
        return []

    # Determinar janela de 6 meses
    actual_max = df[date_col].max()
    actual_min = df[date_col].min()
    chart_min = actual_max - pd.DateOffset(months=n_months - 1)
    chart_min = chart_min.replace(day=1)  # Início do mês
    if chart_min < actual_min:
        chart_min = actual_min.replace(day=1)

    df_chart = df[df[date_col] >= chart_min].copy()
    df_chart['year_month'] = df_chart[date_col].dt.to_period('M')

    all_months = pd.period_range(start=chart_min, end=actual_max, freq='M')
    equipments = sorted(df_chart[equip_col].unique())

    if len(equipments) == 0:
        print("  ⚠ Nenhum equipamento encontrado no período.")
        return []

    # Carregar datas de manutenção
    equip_stats = load_equipment_stats()

    print(f"  Gerando gráficos para {len(equipments)} equipamentos, {len(all_months)} meses...")

    chart_paths = []
    charts_per_page = 8

    for page_idx in range(0, len(equipments), charts_per_page):
        page_equips = equipments[page_idx:page_idx + charts_per_page]
        n_charts = len(page_equips)

        fig, axes = plt.subplots(4, 2, figsize=(16, 22))
        axes_flat = axes.flatten()

        for i, equip in enumerate(page_equips):
            ax = axes_flat[i]
            df_eq = df_chart[df_chart[equip_col] == equip]

            # Agregar por mês
            agg_dict = {prod_col: 'sum'}
            if ref_col:
                agg_dict[ref_col] = 'sum'
            if retrab_col:
                agg_dict[retrab_col] = 'sum'

            monthly = df_eq.groupby('year_month').agg(agg_dict)
            monthly = monthly.reindex(all_months, fill_value=0)

            x = np.arange(len(all_months))
            has_retrab = retrab_col and monthly[retrab_col].sum() > 0
            n_bars = 1 + (1 if ref_col else 0) + (1 if has_retrab else 0)
            width = 0.75 / n_bars

            # Barras de produção
            offset = 0
            bars_prod = ax.bar(
                x + offset * width - (n_bars - 1) * width / 2,
                monthly[prod_col].values, width,
                label='Produção', color='#4CAF50', alpha=0.85, zorder=3
            )
            offset += 1

            # Barras de refugo
            if ref_col:
                bars_ref = ax.bar(
                    x + offset * width - (n_bars - 1) * width / 2,
                    monthly[ref_col].values, width,
                    label='Refugo', color='#F44336', alpha=0.85, zorder=3
                )
                offset += 1

            # Barras de retrabalho (só se houver dados)
            if has_retrab:
                ax.bar(
                    x + offset * width - (n_bars - 1) * width / 2,
                    monthly[retrab_col].values, width,
                    label='Retrabalho', color='#FF9800', alpha=0.85, zorder=3
                )

            # Marcadores de manutenção
            if equip in equip_stats:
                for key in ['ultima_troca', 'penultima_troca']:
                    dt_str = equip_stats[equip].get(key)
                    if dt_str:
                        try:
                            maint_dt = pd.Timestamp(dt_str)
                            maint_period = maint_dt.to_period('M')
                            for j, m in enumerate(all_months):
                                if m == maint_period:
                                    ax.axvline(
                                        x=j, color='#2196F3', linestyle='--',
                                        linewidth=2.5, alpha=0.8, zorder=4
                                    )
                                    ylim = ax.get_ylim()
                                    ax.text(
                                        j, ylim[1] * 0.92, 'Manut.',
                                        fontsize=7, color='#2196F3', ha='center',
                                        fontweight='bold', zorder=5
                                    )
                        except Exception:
                            pass

            # Formatação
            ax.set_title(equip, fontsize=12, fontweight='bold')
            ax.set_xticks(x)
            month_labels = [m.strftime('%b/%Y') for m in all_months]
            ax.set_xticklabels(month_labels, rotation=45, fontsize=8, ha='right')
            ax.set_ylabel('Quantidade (peças)', fontsize=9)
            ax.legend(fontsize=7, loc='upper left')
            ax.grid(axis='y', alpha=0.3, zorder=0)
            ax.set_axisbelow(True)
            ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda val, p: f'{int(val):,}'))

        # Ocultar subplots não usados
        for i in range(n_charts, 8):
            axes_flat[i].set_visible(False)

        periodo_str = f"{chart_min.strftime('%m/%Y')} a {actual_max.strftime('%m/%Y')}"
        fig.suptitle(
            f'Análise Mensal por Equipamento - Produção, Refugo e Manutenção\n'
            f'Período: {periodo_str}',
            fontsize=13, fontweight='bold', y=1.0
        )
        fig.tight_layout(rect=[0, 0, 1, 0.95])

        page_num = page_idx // charts_per_page + 1
        chart_path = output_dir / f"monthly_equip_page_{page_num:02d}.png"
        fig.savefig(chart_path, dpi=150, bbox_inches='tight', facecolor='white')
        plt.close(fig)

        chart_paths.append((chart_path, n_charts))
        print(f"    Página {page_num}: {', '.join(page_equips)}")

    print(f"  ✓ {len(chart_paths)} páginas de gráficos mensais geradas")
    return chart_paths


# =============================================================================
# CONSUMO DE MASSA POR EQUIPAMENTO × COMPOSTO
# =============================================================================

def _load_equipment_compound_aggregate(inicio=None, fim=None):
    """
    Lê data_raw.csv e retorna um DataFrame agregado por (equipamento, composto)
    com colunas: equipamento, composto, pecas_produzidas, kg_consumidos.

    Kg consumidos = Qtd. Produzida × Consumo_de_massa_por_100_peças / 100.
    Retorna None caso os dados ou colunas necessárias não existam.
    """
    data_path = Path("data_raw.csv")
    if not data_path.exists():
        return None

    try:
        df = pd.read_csv(data_path)
    except Exception as e:
        print(f"  ⚠ Erro ao carregar data_raw.csv: {e}")
        return None

    # Detectar colunas dinamicamente (suporte a nomes com acentos)
    date_col = next((c for c in df.columns if 'data' in c.lower() and 'produ' in c.lower()), None)
    prod_col = next((c for c in df.columns if 'produzid' in c.lower() and 'acumul' not in c.lower()), None)
    equip_col = next((c for c in df.columns if c.lower() == 'equipamento'), None)
    if equip_col is None:
        equip_col = next((c for c in df.columns if 'recurso' in c.lower()), None)
    comp_col = next((c for c in df.columns if 'massa' in c.lower() and 'composto' in c.lower()
                     and 'consumo' not in c.lower()), None)
    cons_col = next((c for c in df.columns if 'consumo' in c.lower() and 'massa' in c.lower()
                     and 'acumul' not in c.lower()), None)

    if not all([prod_col, equip_col, comp_col, cons_col]):
        print(f"  ⚠ Colunas necessárias não encontradas em data_raw.csv "
              f"(produção={prod_col}, equip={equip_col}, composto={comp_col}, consumo={cons_col})")
        return None

    # Filtrar período
    if date_col:
        df[date_col] = pd.to_datetime(df[date_col], errors='coerce')
        df = df.dropna(subset=[date_col])
        if inicio:
            df = df[df[date_col] >= pd.Timestamp(inicio)]
        if fim:
            df = df[df[date_col] <= pd.Timestamp(fim)]

    # Tipar numéricos e remover linhas incompletas
    df[prod_col] = pd.to_numeric(df[prod_col], errors='coerce')
    df[cons_col] = pd.to_numeric(df[cons_col], errors='coerce')
    df = df.dropna(subset=[equip_col, comp_col, prod_col, cons_col])

    if len(df) == 0:
        return None

    # Kg = peças × kg/100pçs / 100
    df['_kg'] = df[prod_col] * df[cons_col] / 100.0

    grp = df.groupby([equip_col, comp_col], as_index=False).agg(
        pecas_produzidas=(prod_col, 'sum'),
        kg_consumidos=('_kg', 'sum'),
    )
    grp = grp.rename(columns={equip_col: 'equipamento', comp_col: 'composto'})
    grp = grp.sort_values(['equipamento', 'kg_consumidos'], ascending=[True, False])
    return grp


def generate_equipment_compound_charts(inicio=None, fim=None) -> list:
    """
    Gera gráficos de barras horizontais mostrando, para cada equipamento,
    o consumo de massa (Kg) por tipo de composto. Paginado em 8 equipamentos
    por página para facilitar a leitura no relatório.

    Returns:
        Lista de tuplas (caminho_png, n_charts_na_pagina).
    """
    if not HAS_MATPLOTLIB:
        print("  ⚠ Matplotlib não disponível. Gráficos por composto não serão gerados.")
        return []

    agg = _load_equipment_compound_aggregate(inicio=inicio, fim=fim)
    if agg is None or len(agg) == 0:
        print("  ⚠ Sem dados suficientes para gráficos equipamento × composto.")
        return []

    output_dir = Path("eda_plots") / "compound"
    output_dir.mkdir(parents=True, exist_ok=True)

    equipments = sorted(agg['equipamento'].unique())
    print(f"  Gerando gráficos para {len(equipments)} equipamentos...")

    chart_paths = []
    charts_per_page = 8

    for page_idx in range(0, len(equipments), charts_per_page):
        page_equips = equipments[page_idx:page_idx + charts_per_page]
        n_charts = len(page_equips)

        fig, axes = plt.subplots(4, 2, figsize=(16, 22))
        axes_flat = axes.flatten()

        for i, equip in enumerate(page_equips):
            ax = axes_flat[i]
            eq_df = agg[agg['equipamento'] == equip].sort_values('kg_consumidos', ascending=True)

            compostos = eq_df['composto'].astype(str).tolist()
            kgs = eq_df['kg_consumidos'].values
            pecas = eq_df['pecas_produzidas'].values

            y = np.arange(len(compostos))
            bars = ax.barh(y, kgs, color='#1f77b4', alpha=0.85, zorder=3)

            # Rótulos: Kg (peças)
            xmax = kgs.max() if len(kgs) and kgs.max() > 0 else 1
            for bar, kg, pc in zip(bars, kgs, pecas):
                ax.text(
                    bar.get_width() + xmax * 0.01,
                    bar.get_y() + bar.get_height() / 2,
                    f"{kg:,.0f} kg ({int(pc):,} pç)",
                    va='center', fontsize=7, color='#333333'
                )

            ax.set_yticks(y)
            ax.set_yticklabels(compostos, fontsize=8)
            ax.set_title(f"{equip} — Consumo por Composto", fontsize=11, fontweight='bold')
            ax.set_xlabel('Massa consumida (kg)', fontsize=9)
            ax.set_xlim(0, xmax * 1.25 if xmax > 0 else 1)
            ax.grid(axis='x', alpha=0.3, zorder=0)
            ax.set_axisbelow(True)
            ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda val, p: f'{int(val):,}'))

        # Ocultar subplots não usados
        for i in range(n_charts, 8):
            axes_flat[i].set_visible(False)

        fig.suptitle(
            'Consumo de Massa por Equipamento e Composto\n'
            '(kg consumidos + peças produzidas)',
            fontsize=13, fontweight='bold', y=1.0
        )
        fig.tight_layout(rect=[0, 0, 1, 0.95])

        page_num = page_idx // charts_per_page + 1
        chart_path = output_dir / f"compound_equip_page_{page_num:02d}.png"
        fig.savefig(chart_path, dpi=150, bbox_inches='tight', facecolor='white')
        plt.close(fig)

        chart_paths.append((chart_path, n_charts))
        print(f"    Página {page_num}: {', '.join(page_equips)}")

    print(f"  ✓ {len(chart_paths)} páginas de gráficos equipamento × composto geradas")
    return chart_paths


def generate_equipment_compound_heatmap(inicio=None, fim=None) -> str:
    """
    Gera um heatmap único equipamento × composto mostrando Kg consumidos.

    Returns:
        Caminho do PNG gerado, ou None se indisponível.
    """
    if not HAS_MATPLOTLIB:
        return None

    agg = _load_equipment_compound_aggregate(inicio=inicio, fim=fim)
    if agg is None or len(agg) == 0:
        return None

    output_dir = Path("eda_plots") / "compound"
    output_dir.mkdir(parents=True, exist_ok=True)

    pivot = agg.pivot_table(
        index='equipamento', columns='composto',
        values='kg_consumidos', aggfunc='sum', fill_value=0.0
    )
    # Ordenar compostos pela soma total (mais usado à esquerda)
    col_order = pivot.sum(axis=0).sort_values(ascending=False).index
    pivot = pivot[col_order]

    n_equip, n_comp = pivot.shape
    fig_w = max(12, min(0.55 * n_comp + 4, 26))
    fig_h = max(6, min(0.35 * n_equip + 3, 22))

    fig, ax = plt.subplots(figsize=(fig_w, fig_h))
    data = pivot.values
    vmax = data.max() if data.size and data.max() > 0 else 1
    im = ax.imshow(data, aspect='auto', cmap='YlOrRd', vmin=0, vmax=vmax)

    ax.set_xticks(np.arange(n_comp))
    ax.set_xticklabels(pivot.columns, rotation=45, ha='right', fontsize=8)
    ax.set_yticks(np.arange(n_equip))
    ax.set_yticklabels(pivot.index, fontsize=8)
    ax.set_xlabel('Composto', fontsize=10)
    ax.set_ylabel('Equipamento', fontsize=10)
    ax.set_title('Heatmap — Consumo de Massa (kg) por Equipamento × Composto',
                 fontsize=12, fontweight='bold')

    # Anotar células com valor > 0 (em milhares de kg se grande)
    for i in range(n_equip):
        for j in range(n_comp):
            val = data[i, j]
            if val > 0:
                text = f"{val/1000:.1f}k" if val >= 1000 else f"{val:.0f}"
                color = 'white' if val > vmax * 0.55 else 'black'
                ax.text(j, i, text, ha='center', va='center', fontsize=6, color=color)

    cbar = fig.colorbar(im, ax=ax, shrink=0.8)
    cbar.set_label('Massa consumida (kg)', fontsize=9)

    fig.tight_layout()
    out_path = output_dir / "compound_equipamento_heatmap.png"
    fig.savefig(out_path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    print(f"  ✓ Heatmap equipamento × composto: {out_path}")
    return str(out_path)


def build_equipment_compound_summary(inicio=None, fim=None) -> list:
    """
    Constrói a tabela resumo (equipamento, composto, peças, Kg) usada na
    seção 17 do relatório.

    Returns:
        Lista de linhas como listas de strings, pronta para ser alimentada
        em um reportlab.platypus.Table. Cabeçalho é a primeira linha.
        Retorna None quando não há dados.
    """
    agg = _load_equipment_compound_aggregate(inicio=inicio, fim=fim)
    if agg is None or len(agg) == 0:
        return None

    rows = [["Equipamento", "Composto", "Peças produzidas", "Massa consumida (kg)"]]
    for _, r in agg.iterrows():
        rows.append([
            str(r['equipamento']),
            str(r['composto']),
            f"{int(r['pecas_produzidas']):,}".replace(",", "."),
            f"{r['kg_consumidos']:,.1f}".replace(",", "X").replace(".", ",").replace("X", "."),
        ])
    return rows


def main(inicio=None, fim=None, suffix="", version=None, **kwargs) -> dict:
    """
    Função principal - Etapa 6: Geração de Relatório.

    Args:
        inicio: Data de início do período (YYYY-MM-DD), opcional
        fim: Data de fim do período (YYYY-MM-DD), opcional
        suffix: Sufixo para o nome do relatório (ex: "_v1"), opcional
        version: Versão fixa do relatório (ex: "R12_v1"), sobrescreve auto-incremento

    Returns:
        Dicionário com resultados da execução
    """
    print("=" * 60)
    print("ETAPA 6: GERAÇÃO DE RELATÓRIO")
    print("=" * 60)

    # Verificar arquivos necessários
    required_files = ["best_model.joblib"]
    missing = [f for f in required_files if not Path(f).exists()]

    if missing:
        print(f"\n⚠ Arquivos não encontrados: {', '.join(missing)}")
        print("Execute as etapas anteriores do pipeline primeiro.")
        return {"status": "error", "message": "Missing required files"}

    # Carregar resultados
    print("\n[1/4] Carregando resultados do pipeline...")
    results = load_pipeline_results()

    # Gerar gráficos mensais por equipamento
    print("\n[2/4] Gerando gráficos mensais por equipamento...")
    monthly_charts = generate_monthly_equipment_charts(inicio=inicio, fim=fim)
    results["monthly_charts"] = monthly_charts

    # Gerar gráficos e tabela de consumo por equipamento × composto
    print("\n[2b/4] Gerando análise de consumo por equipamento × composto...")
    results["compound_charts"] = generate_equipment_compound_charts(inicio=inicio, fim=fim)
    results["compound_heatmap"] = generate_equipment_compound_heatmap(inicio=inicio, fim=fim)
    results["compound_summary_rows"] = build_equipment_compound_summary(inicio=inicio, fim=fim)

    # Determinar versão
    if not version:
        version = get_report_version(suffix=suffix or "")

    # Gerar relatório
    print(f"\n[3/4] Gerando relatório versão {version}...")

    if HAS_REPORTLAB:
        output_path = f"Relatorio_SABO_{version}.pdf"
        generated_path = generate_pdf_report(results, output_path, inicio=inicio, fim=fim)
        print(f"  ✓ Relatório PDF gerado: {generated_path}")
    else:
        output_path = f"Relatorio_SABO_{version}.txt"
        generated_path = generate_text_report(results, output_path)
        print(f"  ✓ Relatório TXT gerado: {generated_path}")
        print("  ⚠ Instale 'reportlab' para gerar PDF: pip install reportlab")

    # Resumo
    print("\n[4/4] Finalizando...")

    print("\n" + "=" * 60)
    print("ETAPA 6 CONCLUÍDA")
    print("=" * 60)
    print(f"\nRelatório gerado: {generated_path}")

    # Informações do modelo
    model_name = results.get("model", {}).get("name", "N/A")
    metrics = results.get("model", {}).get("metrics", {})

    print(f"\nModelo documentado: {model_name.upper()}")
    if metrics:
        print(f"  R²:  {metrics.get('r2', 0):.4f}")
        print(f"  MSE: {metrics.get('mse', 0):.2f}")
        print(f"  MAE: {metrics.get('mae', 0):.2f}")

    return {
        "status": "success",
        "output_file": generated_path,
        "version": version,
        "model": model_name,
        "best_metrics": metrics,
    }


if __name__ == "__main__":
    main()
